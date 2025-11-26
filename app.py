import streamlit as st
import pandas as pd
import duckdb
import requests
import json
import re
from datetime import datetime
import io
import base64
from typing import Dict, List, Tuple
import pyarrow as pa
import pyarrow.parquet as pq
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import asyncio
import aiohttp
from asyncio import Semaphore
import time
import hashlib
import pickle
from pathlib import Path
import os
from collections import deque

# Page config
st.set_page_config(
    page_title="QA Coaching Intelligence",
    page_icon="📈",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Custom CSS
st.markdown("""
<style>
    @import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;600;700&display=swap');
    
    * {
        font-family: 'Inter', sans-serif;
    }
    
    .main {
        background: #ffffff;
        background-attachment: fixed;
    }
    
    .stApp {
        background: transparent;
    }
    
    div[data-testid="stMetricValue"] {
        font-size: 2.5rem;
        font-weight: 700;
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
    }
    
    .stTabs [data-baseweb="tab-list"] {
        gap: 24px;
        background: rgba(255,255,255,0.95);
        padding: 20px;
        border-radius: 15px;
        backdrop-filter: blur(10px);
    }
    
    .stTabs [data-baseweb="tab"] {
        height: 60px;
        padding: 0 30px;
        font-weight: 600;
        font-size: 1.1rem;
        border-radius: 10px;
        transition: all 0.3s ease;
    }
    
    .stTabs [data-baseweb="tab"]:hover {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        color: white;
    }
    
    .stTabs [aria-selected="true"] {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        color: white !important;
    }
    
    .upload-box {
        background: rgba(255,255,255,0.95);
        padding: 40px;
        border-radius: 20px;
        text-align: center;
        border: 3px dashed #667eea;
        transition: all 0.3s ease;
    }
    
    .upload-box:hover {
        border-color: #764ba2;
        transform: translateY(-5px);
        box-shadow: 0 20px 40px rgba(102,126,234,0.3);
    }
    
    .metric-card {
        background: white;
        border: 1px solid #e2e8f0;
        padding: 24px;
        border-radius: 12px;
        box-shadow: 0 1px 3px rgba(0,0,0,0.05);
        transition: all 0.3s ease;
    }
    
    .metric-card:hover {
        transform: translateY(-2px);
        box-shadow: 0 4px 12px rgba(0,0,0,0.1);
        border-color: #0ea5e9;
    }
    
    .agent-card {
        background: white;
        padding: 24px;
        border-radius: 12px;
        margin: 15px 0;
        box-shadow: 0 1px 3px rgba(0,0,0,0.05);
        border: 1px solid #e2e8f0;
        border-left: 4px solid #0ea5e9;
        transition: all 0.3s ease;
    }
    
    .agent-card:hover {
        box-shadow: 0 4px 12px rgba(0,0,0,0.1);
        border-left-color: #0284c7;
    }
    
    .theme-badge {
        display: inline-block;
        padding: 6px 14px;
        border-radius: 16px;
        font-weight: 600;
        font-size: 0.85rem;
        margin: 5px;
    }
    
    .priority-high {
        background: #fef2f2;
        color: #dc2626;
        border: 1px solid #fecaca;
    }
    
    .priority-medium {
        background: #fffbeb;
        color: #d97706;
        border: 1px solid #fde68a;
    }
    
    .priority-low {
        background: #f0fdf4;
        color: #16a34a;
        border: 1px solid #bbf7d0;
    }
    
    .chat-message {
        padding: 15px 20px;
        border-radius: 15px;
        margin: 10px 0;
        animation: fadeIn 0.3s ease;
    }
    
    @keyframes fadeIn {
        from { opacity: 0; transform: translateY(10px); }
        to { opacity: 1; transform: translateY(0); }
    }
    
    .user-message {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        color: white;
        margin-left: 20%;
    }
    
    .assistant-message {
        background: rgba(255,255,255,0.95);
        color: #333;
        margin-right: 20%;
        border: 2px solid #667eea;
    }
    
    .stButton>button {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        color: white;
        border: none;
        padding: 15px 40px;
        font-size: 1.1rem;
        font-weight: 600;
        border-radius: 10px;
        transition: all 0.3s ease;
    }
    
    .stButton>button:hover {
        transform: translateY(-3px);
        box-shadow: 0 10px 25px rgba(102,126,234,0.4);
    }
    
    .stSelectbox, .stMultiSelect {
        background: white;
        border-radius: 10px;
    }
    
    .success-banner {
        background: linear-gradient(135deg, #11998e 0%, #38ef7d 100%);
        color: white;
        padding: 20px;
        border-radius: 15px;
        font-weight: 600;
        text-align: center;
        animation: slideDown 0.5s ease;
    }
    
    @keyframes slideDown {
        from { transform: translateY(-20px); opacity: 0; }
        to { transform: translateY(0); opacity: 1; }
    }
</style>
""", unsafe_allow_html=True)

# Initialize session state
if 'duckdb_conn' not in st.session_state:
    st.session_state.duckdb_conn = duckdb.connect(':memory:')
if 'coaching_insights' not in st.session_state:
    st.session_state.coaching_insights = {}
if 'chat_history' not in st.session_state:
    st.session_state.chat_history = []
if 'processed' not in st.session_state:
    st.session_state.processed = False
if 'processing_stats' not in st.session_state:
    st.session_state.processing_stats = {
        'total_batches': 0,
        'completed_batches': 0,
        'failed_batches': 0,
        'start_time': None,
        'total_tokens': 0
    }
if 'data_loaded' not in st.session_state:
    st.session_state.data_loaded = False
if 'pre_analysis_done' not in st.session_state:
    st.session_state.pre_analysis_done = False

# Model configurations
MODELS = {
    "deepseek/deepseek-chat:free": {
        "name": "DeepSeek Chat",
        "rating": "5/5",
        "best_for": "Advanced reasoning & analysis",
        "speed": "Fast",
        "recommended": True
    },
    "deepseek/deepseek-r1-distill-llama-70b:free": {
        "name": "DeepSeek R1 Distill",
        "rating": "5/5",
        "best_for": "Deep reasoning on complex cases",
        "speed": "Medium"
    },
    "meta-llama/llama-3.3-70b-instruct:free": {
        "name": "Llama 3.3 70B",
        "rating": "5/5",
        "best_for": "Balanced performance & quality",
        "speed": "Medium"
    },
    "x-ai/grok-4.1-fast:free": {
        "name": "Grok 4.1 Fast",
        "rating": "4/5",
        "best_for": "Fast reasoning & insights",
        "speed": "Very Fast"
    },
    "google/gemma-3-27b-it:free": {
        "name": "Gemma 3 27B",
        "rating": "4/5",
        "best_for": "Google's efficient model",
        "speed": "Fast"
    },
    "qwen/qwen-2.5-72b-instruct:free": {
        "name": "Qwen 2.5 72B",
        "rating": "4/5",
        "best_for": "Structured analysis",
        "speed": "Fast"
    },
    "openai/gpt-oss-20b:free": {
        "name": "GPT OSS 20B",
        "rating": "3/5",
        "best_for": "Basic coaching themes",
        "speed": "Fast"
    },
    "meituan/longcat-flash-chat:free": {
        "name": "LongCat Flash",
        "rating": "3/5",
        "best_for": "Quick chat analysis",
        "speed": "Very Fast"
    },
    "microsoft/mai-ds-r1:free": {
        "name": "MAI DS R1",
        "rating": "4/5",
        "best_for": "Microsoft's reasoning model",
        "speed": "Medium"
    },
    "mistralai/mistral-7b-instruct:free": {
        "name": "Mistral 7B",
        "rating": "3/5",
        "best_for": "Quick basic analysis",
        "speed": "Very Fast"
    },
    "gryphe/mythomax-l2-13b:free": {
        "name": "MythoMax L2 13B",
        "rating": "3/5",
        "best_for": "Creative coaching suggestions",
        "speed": "Fast"
    },
    "mistralai/mistral-nemo:free": {
        "name": "Mistral Nemo",
        "rating": "3/5",
        "best_for": "Fast Q&A chat",
        "speed": "Very Fast"
    }
}

# Default coaching themes with 5C framework mapping
DEFAULT_THEMES = [
    # Connection (Building Rapport)
    "Active Listening & Acknowledgment",
    "Empathy & Emotional Intelligence",
    "Building Rapport & Trust",
    "Professional Tone & Language",
    "Personalization & Context Awareness",
    "Emotional Regulation",
    
    # Clarity (Clear Communication)
    "Clear Communication & Articulation",
    "Process Adherence & Documentation",
    "Product Knowledge & Accuracy",
    "Jargon-Free Communication",
    "Step-by-Step Guidance",
    "Confirmation & Recap Skills",
    
    # Commitment (Ownership & Follow-through)
    "First Call Resolution",
    "Solution Offering & Alternatives",
    "Follow-up & Closure Quality",
    "Ownership Language Usage",
    "Timeline Setting & Commitments",
    "Accountability & Promises",
    
    # Challenge (Problem-solving)
    "Problem Diagnosis & Root Cause",
    "Handling Difficult Customers",
    "Escalation Judgment & Timing",
    "De-escalation Techniques",
    "Objection Handling",
    "Critical Thinking & Analysis",
    
    # Change (Adaptability & Growth)
    "Response Time Management",
    "Proactive Communication",
    "Managing Customer Expectations",
    "Adaptability & Flexibility",
    "Feedback Responsiveness",
    "Continuous Improvement Mindset"
]

# 5C Framework mapping
THEME_TO_5C = {
    "Connection": [
        "Active Listening & Acknowledgment",
        "Empathy & Emotional Intelligence",
        "Building Rapport & Trust",
        "Professional Tone & Language",
        "Personalization & Context Awareness",
        "Emotional Regulation"
    ],
    "Clarity": [
        "Clear Communication & Articulation",
        "Process Adherence & Documentation",
        "Product Knowledge & Accuracy",
        "Jargon-Free Communication",
        "Step-by-Step Guidance",
        "Confirmation & Recap Skills"
    ],
    "Commitment": [
        "First Call Resolution",
        "Solution Offering & Alternatives",
        "Follow-up & Closure Quality",
        "Ownership Language Usage",
        "Timeline Setting & Commitments",
        "Accountability & Promises"
    ],
    "Challenge": [
        "Problem Diagnosis & Root Cause",
        "Handling Difficult Customers",
        "Escalation Judgment & Timing",
        "De-escalation Techniques",
        "Objection Handling",
        "Critical Thinking & Analysis"
    ],
    "Change": [
        "Response Time Management",
        "Proactive Communication",
        "Managing Customer Expectations",
        "Adaptability & Flexibility",
        "Feedback Responsiveness",
        "Continuous Improvement Mindset"
    ]
}

# 5C Icons - SVG paths for each pillar
FIVE_C_ICONS = {
    "Connection": '<svg width="28" height="28" viewBox="0 0 24 24" fill="none" stroke="#3b82f6" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"><path d="M17 21v-2a4 4 0 0 0-4-4H5a4 4 0 0 0-4 4v2"></path><circle cx="9" cy="7" r="4"></circle><path d="M23 21v-2a4 4 0 0 0-3-3.87"></path><path d="M16 3.13a4 4 0 0 1 0 7.75"></path></svg>',
    "Clarity": '<svg width="28" height="28" viewBox="0 0 24 24" fill="none" stroke="#8b5cf6" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"><circle cx="12" cy="12" r="10"></circle><line x1="12" y1="16" x2="12" y2="12"></line><line x1="12" y1="8" x2="12.01" y2="8"></line></svg>',
    "Commitment": '<svg width="28" height="28" viewBox="0 0 24 24" fill="none" stroke="#0ea5e9" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"><polyline points="9 11 12 14 22 4"></polyline><path d="M21 12v7a2 2 0 0 1-2 2H5a2 2 0 0 1-2-2V5a2 2 0 0 1 2-2h11"></path></svg>',
    "Challenge": '<svg width="28" height="28" viewBox="0 0 24 24" fill="none" stroke="#10b981" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"><path d="M22 12h-4l-3 9L9 3l-3 9H2"></path></svg>',
    "Change": '<svg width="28" height="28" viewBox="0 0 24 24" fill="none" stroke="#f59e0b" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"><polyline points="23 4 23 10 17 10"></polyline><path d="M20.49 15a9 9 0 1 1-2.12-9.36L23 10"></path></svg>'
}

# 5C Colors
FIVE_C_COLORS = {
    "Connection": "#3b82f6",  # Blue
    "Clarity": "#8b5cf6",     # Purple  
    "Commitment": "#0ea5e9",  # Sky blue
    "Challenge": "#10b981",   # Green
    "Change": "#f59e0b"       # Amber
}

class RateLimiter:
    """Token bucket rate limiter for API calls"""
    def __init__(self, calls_per_minute: int):
        self.calls_per_minute = calls_per_minute
        self.tokens = calls_per_minute
        self.last_update = time.time()
        self.lock = asyncio.Lock()
    
    async def acquire(self):
        async with self.lock:
            now = time.time()
            elapsed = now - self.last_update
            
            # Refill tokens based on elapsed time
            self.tokens = min(
                self.calls_per_minute,
                self.tokens + (elapsed * self.calls_per_minute / 60)
            )
            self.last_update = now
            
            if self.tokens >= 1:
                self.tokens -= 1
                return
            
            # Wait until next token available
            wait_time = (1 - self.tokens) * 60 / self.calls_per_minute
            await asyncio.sleep(wait_time)
            self.tokens = 0
            self.last_update = time.time()

def parse_transcript_chunk(chunk_data):
    """Parse a chunk of transcripts - used for multiprocessing"""
    results = []
    for row_data in chunk_data:
        call_id, agent_name, transcript_text, sentiment = row_data
        turns = parse_multiline_transcript(str(transcript_text))
        
        for turn in turns:
            results.append({
                'call_id': call_id,
                'agent': agent_name,
                'timestamp': turn['timestamp'],
                'speaker': turn['speaker'],
                'message': turn['message'],
                'sentiment_score': sentiment,
                'original_transcript': transcript_text
            })
    
    return results

def parse_transcripts_parallel(df, call_id_col, agent_col, transcript_col, sentiment_col, num_workers=None):
    """Parse transcripts in parallel using concurrent.futures (Streamlit Cloud compatible)"""
    from concurrent.futures import ProcessPoolExecutor, as_completed
    from multiprocessing import cpu_count
    import os
    
    if num_workers is None:
        num_workers = max(1, cpu_count() - 1)
    
    # Prepare data
    chunk_data = []
    for idx, row in df.iterrows():
        call_id = row[call_id_col]
        agent_name = row[agent_col]
        transcript_text = row[transcript_col]
        sentiment = None
        if sentiment_col and sentiment_col != "None":
            sentiment = row.get(sentiment_col)
        
        chunk_data.append((call_id, agent_name, transcript_text, sentiment))
    
    # Split into chunks
    chunk_size = max(50, len(chunk_data) // (num_workers * 8))  # Smaller chunks for responsiveness
    chunks = [chunk_data[i:i + chunk_size] for i in range(0, len(chunk_data), chunk_size)]
    
    # Process in parallel with progress updates
    expanded_rows = []
    
    # Use ProcessPoolExecutor with spawn context for Streamlit Cloud
    with ProcessPoolExecutor(max_workers=num_workers) as executor:
        # Submit all chunks
        future_to_chunk = {executor.submit(parse_transcript_chunk, chunk): i for i, chunk in enumerate(chunks)}
        
        # Process results as they complete (allows yielding control)
        for future in as_completed(future_to_chunk):
            try:
                chunk_result = future.result(timeout=30)  # 30s timeout per chunk
                expanded_rows.extend(chunk_result)
            except Exception as e:
                print(f"Chunk processing error: {str(e)}")
                continue
    
    return expanded_rows

def redact_pii(text: str) -> str:
    """Redact PII from text"""
    # Email
    text = re.sub(r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b', '[EMAIL_REDACTED]', text)
    # Phone
    text = re.sub(r'\b\d{3}[-.]?\d{3}[-.]?\d{4}\b', '[PHONE_REDACTED]', text)
    # SSN
    text = re.sub(r'\b\d{3}-\d{2}-\d{4}\b', '[SSN_REDACTED]', text)
    # Credit card (basic)
    text = re.sub(r'\b\d{4}[\s-]?\d{4}[\s-]?\d{4}[\s-]?\d{4}\b', '[CARD_REDACTED]', text)
    return text

def normalize_speaker(speaker: str) -> str:
    """Normalize speaker labels"""
    speaker_lower = speaker.lower()
    if any(x in speaker_lower for x in ['agent', 'representative', 'rep']):
        return 'agent'
    elif any(x in speaker_lower for x in ['customer', 'consumer', 'client']):
        return 'customer'
    return speaker_lower

# Compile regex patterns once for performance
BRACKET_PATTERN = re.compile(r'\[([\d:]+)\s+([^\]]+)\]:\s*\n?\s*(.*?)(?=\[[\d:]+\s+[^\]]+\]:|$)', re.DOTALL)
PIPE_PATTERN = re.compile(r'(\d{4}-\d{2}-\d{2}\s+\d{2}:\d{2}:\d{2}\s+[+-]\d{4})\s+([^:]+):\s*(.*)')
HTML_BR_PATTERN = re.compile(r'<br\s*/?>', re.IGNORECASE)
HTML_TAG_PATTERN = re.compile(r'<[^>]+>')

def parse_multiline_transcript(transcript_text: str) -> List[Dict]:
    """Parse multiline transcript from single cell into conversation turns
    
    Handles formats:
    1. Bracket with newline: "[12:30:08 AGENT]:\n message"
    2. Bracket inline: "[12:30:08 AGENT]: message"
    3. Pipe-separated: "2025-02-07 13:17:57 +0000 Consumer: Hi! | 2025-02-07 13:18:01 +0000 Agent: Hello"
    """
    if not transcript_text or not isinstance(transcript_text, str):
        return []
    
    turns = []
    
    # Check if pipe-separated format
    if '|' in transcript_text:
        segments = transcript_text.split('|')
        for segment in segments:
            segment = segment.strip()
            if not segment:
                continue
            match = PIPE_PATTERN.match(segment)
            if match:
                timestamp, speaker, message = match.groups()
                # Clean HTML tags
                message = HTML_BR_PATTERN.sub(' ', message)
                message = HTML_TAG_PATTERN.sub('', message)
                message = ' '.join(message.split())
                if message:
                    turns.append({
                        'timestamp': timestamp,
                        'speaker': normalize_speaker(speaker.strip()),
                        'message': redact_pii(message.strip())
                    })
    else:
        # Try bracket format
        matches = BRACKET_PATTERN.findall(transcript_text)
        for match in matches:
            timestamp, speaker, message = match
            # Clean HTML tags
            message = HTML_BR_PATTERN.sub(' ', message)
            message = HTML_TAG_PATTERN.sub('', message)
            message = ' '.join(message.split())
            if message:
                turns.append({
                    'timestamp': timestamp.strip(),
                    'speaker': normalize_speaker(speaker.strip()),
                    'message': redact_pii(message.strip())
                })
    
    return turns

def convert_to_parquet(df: pd.DataFrame, filename: str) -> bytes:
    """Convert dataframe to parquet bytes"""
    table = pa.Table.from_pandas(df)
    buf = io.BytesIO()
    pq.write_table(table, buf, compression='snappy')
    buf.seek(0)
    return buf.getvalue()

def load_file_to_dataframe(uploaded_file) -> pd.DataFrame:
    """Load various file formats to raw dataframe (no parsing)"""
    file_ext = uploaded_file.name.split('.')[-1].lower()
    
    if file_ext == 'csv':
        return pd.read_csv(uploaded_file)
    elif file_ext in ['xlsx', 'xls']:
        return pd.read_excel(uploaded_file)
    elif file_ext == 'parquet':
        return pd.read_parquet(uploaded_file)
    elif file_ext == 'txt':
        # For TXT, create simple dataframe
        content = uploaded_file.read().decode('utf-8')
        return pd.DataFrame([{
            'call_id': 'CALL_0001',
            'agent': 'Unknown',
            'transcript': content
        }])
    else:
        st.error(f"Unsupported file format: {file_ext}")
        return None

async def call_llm_async(
    session: aiohttp.ClientSession,
    model: str,
    messages: List[Dict],
    temperature: float = 0.3,
    is_json: bool = True,
    provider: str = "openrouter",
    api_key: str = None,
    local_url: str = None,
    rate_limiter: RateLimiter = None
) -> Dict:
    """Async LLM call supporting OpenRouter and Local LLM"""
    
    if rate_limiter:
        await rate_limiter.acquire()
    
    payload = {
        "model": model,
        "messages": messages,
        "temperature": temperature,
        "max_tokens": 2000
    }
    
    if is_json and provider == "openrouter":
        payload["response_format"] = {"type": "json_object"}
    
    try:
        if provider == "openrouter":
            url = "https://openrouter.ai/api/v1/chat/completions"
            headers = {
                "Content-Type": "application/json",
                "Authorization": f"Bearer {api_key}" if api_key else ""
            }
        else:  # local LLM
            url = local_url or "http://localhost:1234/v1/chat/completions"
            headers = {"Content-Type": "application/json"}
        
        async with session.post(url, json=payload, headers=headers, timeout=aiohttp.ClientTimeout(total=60)) as response:
            if response.status == 200:
                return await response.json()
            else:
                error_text = await response.text()
                return {"error": f"HTTP {response.status}: {error_text}"}
                
    except asyncio.TimeoutError:
        return {"error": "Request timeout"}
    except Exception as e:
        return {"error": str(e)}

def call_llm(model: str, messages: List[Dict], temperature: float = 0.3, is_json: bool = True) -> Dict:
    """Synchronous wrapper for backward compatibility (used in chat)"""
    provider = st.session_state.get('llm_provider', 'openrouter')
    api_key = st.session_state.get('openrouter_api_key')
    local_url = st.session_state.get('local_llm_url')
    
    if provider == "openrouter":
        url = "https://openrouter.ai/api/v1/chat/completions"
        headers = {
            "Content-Type": "application/json",
            "Authorization": f"Bearer {api_key}" if api_key else ""
        }
    else:
        url = local_url or "http://localhost:1234/v1/chat/completions"
        headers = {"Content-Type": "application/json"}
    
    payload = {
        "model": model,
        "messages": messages,
        "temperature": temperature,
        "max_tokens": 2000
    }
    
    if is_json and provider == "openrouter":
        payload["response_format"] = {"type": "json_object"}
    
    try:
        response = requests.post(url, headers=headers, json=payload, timeout=60)
        response.raise_for_status()
        return response.json()
    except Exception as e:
        st.error(f"LLM API Error: {str(e)}")
        return None

def process_agent_batch(agent_name: str, calls_df: pd.DataFrame, themes: List[str], model: str) -> Dict:
    """Process batch of calls for an agent (kept for backward compatibility)"""
    
    # Create compressed context
    call_summaries = []
    for idx, call_id in enumerate(calls_df['call_id'].unique()[:10], 1):
        call_data = calls_df[calls_df['call_id'] == call_id]
        
        # Get sentiment if available
        sentiment = call_data['sentiment_score'].mean() if 'sentiment_score' in call_data.columns else None
        
        # Compress conversation
        agent_msgs = call_data[call_data['speaker'] == 'agent']['message'].tolist()
        customer_msgs = call_data[call_data['speaker'] == 'customer']['message'].tolist()
        
        summary = f"Call {idx}:"
        if sentiment:
            summary += f" [Sentiment: {sentiment:.2f}]"
        summary += f"\n- Customer issues: {' | '.join(customer_msgs[:3])}"
        summary += f"\n- Agent responses: {' | '.join(agent_msgs[:3])}"
        
        call_summaries.append(summary)
    
    # Build prompt
    system_prompt = f"""You are an expert contact center QA analyst. Analyze agent performance and identify coaching opportunities.

Focus on these themes: {', '.join(themes)}

Provide response as JSON with this exact structure:
{{
    "agent": "agent_name",
    "calls_analyzed": number,
    "coaching_themes": [
        {{
            "theme": "theme name from provided list",
            "priority": "high|medium|low",
            "frequency": number,
            "examples": ["specific example 1", "specific example 2"],
            "recommendation": "specific actionable advice"
        }}
    ],
    "strengths": ["strength 1", "strength 2"],
    "overall_sentiment_correlation": "insight about sentiment and performance"
}}

CRITICAL: Only use themes from the provided list. Be specific and data-driven."""

    user_prompt = f"""Agent: {agent_name}
Calls analyzed: {len(calls_df['call_id'].unique())}

Call Summaries:
{chr(10).join(call_summaries)}

Identify top 3-5 coaching opportunities with specific examples and actionable recommendations."""

    messages = [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": user_prompt}
    ]
    
    response = call_llm(model, messages)
    
    if response and 'choices' in response:
        try:
            content = response['choices'][0]['message']['content']
            # Handle markdown code blocks
            content = re.sub(r'```json\n?', '', content)
            content = re.sub(r'```\n?', '', content)
            return json.loads(content.strip())
        except json.JSONDecodeError as e:
            st.error(f"Failed to parse LLM response for {agent_name}: {str(e)}")
            return None
    
    return None

async def process_agent_batch_async(
    session: aiohttp.ClientSession,
    agent_name: str,
    calls_df: pd.DataFrame,
    themes: List[str],
    model: str,
    provider: str,
    api_key: str,
    local_url: str,
    rate_limiter: RateLimiter,
    semaphore: Semaphore
) -> Tuple[str, Dict]:
    """Async process batch of calls for an agent"""
    
    async with semaphore:
        try:
            # Get unique call IDs for this agent (already filtered to 4-5 calls by caller)
            unique_calls = calls_df['call_id'].unique()
            
            # Create compressed context
            call_summaries = []
            for idx, call_id in enumerate(unique_calls, 1):
                call_data = calls_df[calls_df['call_id'] == call_id]
                
                sentiment = call_data['sentiment_score'].mean() if 'sentiment_score' in call_data.columns else None
                
                # Get only first 3 messages from each side to keep it concise
                agent_msgs = call_data[call_data['speaker'] == 'agent']['message'].tolist()[:3]
                customer_msgs = call_data[call_data['speaker'] == 'customer']['message'].tolist()[:3]
                
                # Compress messages further (first 100 chars each)
                agent_msgs_short = [msg[:100] + '...' if len(msg) > 100 else msg for msg in agent_msgs]
                customer_msgs_short = [msg[:100] + '...' if len(msg) > 100 else msg for msg in customer_msgs]
                
                summary = f"Call {idx}:"
                if sentiment:
                    summary += f" [Sentiment: {sentiment:.2f}]"
                summary += f"\n- Customer: {' | '.join(customer_msgs_short)}"
                summary += f"\n- Agent: {' | '.join(agent_msgs_short)}"
                
                call_summaries.append(summary)
            
            system_prompt = f"""You are an expert contact center QA analyst. Analyze agent performance and identify coaching opportunities.

Focus on these themes: {', '.join(themes[:10])}

Provide response as JSON with this exact structure:
{{
    "agent": "{agent_name}",
    "calls_analyzed": {len(unique_calls)},
    "coaching_themes": [
        {{
            "theme": "theme name from provided list",
            "priority": "high|medium|low",
            "frequency": 1,
            "examples": ["brief example"],
            "recommendation": "specific actionable advice"
        }}
    ],
    "strengths": ["strength 1", "strength 2"],
    "overall_sentiment_correlation": "brief insight"
}}

CRITICAL: Only use themes from the provided list. Be specific and data-driven. Keep response concise."""

            user_prompt = f"""Agent: {agent_name}
Calls analyzed: {len(unique_calls)}

Call Summaries:
{chr(10).join(call_summaries)}

Identify top 3 coaching opportunities with specific examples and actionable recommendations."""

            messages = [
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ]
            
            response = await call_llm_async(
                session, model, messages,
                provider=provider,
                api_key=api_key,
                local_url=local_url,
                rate_limiter=rate_limiter
            )
            
            if not response:
                print(f"[ERROR] {agent_name}: No response from LLM")
                return (agent_name, None)
            
            if 'error' in response:
                print(f"[ERROR] {agent_name}: {response['error']}")
                return (agent_name, None)
            
            if 'choices' not in response:
                print(f"[ERROR] {agent_name}: No 'choices' in response: {response}")
                return (agent_name, None)
            
            try:
                content = response['choices'][0]['message']['content']
                content = re.sub(r'```json\n?', '', content)
                content = re.sub(r'```\n?', '', content)
                result = json.loads(content.strip())
                print(f"[SUCCESS] {agent_name}: Generated {len(result.get('coaching_themes', []))} themes")
                return (agent_name, result)
            except json.JSONDecodeError as e:
                print(f"[ERROR] {agent_name}: JSON parse failed: {str(e)}")
                print(f"[ERROR] Content: {content[:200]}")
                return (agent_name, None)
        
        except Exception as e:
            print(f"[ERROR] {agent_name}: Exception: {str(e)}")
            import traceback
            traceback.print_exc()
            return (agent_name, None)

async def process_all_agents_parallel(
    agents_data: List[Tuple[str, pd.DataFrame]],
    themes: List[str],
    model: str,
    provider: str,
    api_key: str,
    local_url: str,
    max_concurrent: int = 10,
    calls_per_minute: int = 50
) -> Dict:
    """Process all agents in parallel with rate limiting"""
    
    rate_limiter = RateLimiter(calls_per_minute)
    semaphore = Semaphore(max_concurrent)
    
    async with aiohttp.ClientSession() as session:
        tasks = [
            process_agent_batch_async(
                session, agent_name, agent_df, themes, model,
                provider, api_key, local_url, rate_limiter, semaphore
            )
            for agent_name, agent_df in agents_data
        ]
        
        results = await asyncio.gather(*tasks, return_exceptions=True)
    
    insights = {}
    for result in results:
        if isinstance(result, tuple) and result[1]:
            agent_name, data = result
            insights[agent_name] = data
    
    return insights

def generate_email_share_link(agent_name: str, agent_data: Dict) -> str:
    """Generate mailto link with plain text coaching plan (renders in all email clients)"""
    import urllib.parse
    
    themes = agent_data.get('coaching_themes', [])
    calls = agent_data.get('calls_analyzed', 0)
    strengths = agent_data.get('strengths', [])
    
    # Build plain text email body that will render nicely
    text_body = f"""COACHING PLAN FOR {agent_name.upper()}
{'=' * 60}

SUMMARY
   • Calls Analyzed: {calls}
   • Date Generated: {datetime.now().strftime('%B %d, %Y')}

{'=' * 60}

COACHING THEMES
"""
    
    for idx, theme in enumerate(themes, 1):
        theme_name = theme.get('theme', '')
        priority = theme.get('priority', 'low')
        recommendation = theme.get('recommendation', '')
        examples = theme.get('examples', [])
        
        priority_label = {
            'high': '[HIGH PRIORITY]',
            'medium': '[MEDIUM PRIORITY]',
            'low': '[LOW PRIORITY]'
        }.get(priority, '[LOW]')
        
        text_body += f"""
{idx}. {theme_name.upper()}
   Priority: {priority_label}
   
   Recommendation:
   {recommendation}
"""
        
        if examples:
            text_body += f"""   
   Examples:
"""
            for ex in examples[:2]:
                text_body += f"   • {ex}\n"
        
        text_body += "\n" + ("-" * 60) + "\n"
    
    # Strengths
    if strengths:
        text_body += f"""
STRENGTHS
"""
        for strength in strengths:
            text_body += f"   • {strength}\n"
        
        text_body += "\n" + ("=" * 60) + "\n"
    
    text_body += """
Generated by QA Coaching Intelligence Platform
"""
    
    # URL encode for mailto
    subject = f"Coaching Plan - {agent_name}"
    body_encoded = urllib.parse.quote(text_body)
    
    mailto_link = f"mailto:?subject={urllib.parse.quote(subject)}&body={body_encoded}"
    
    return mailto_link

# Session Management Functions
def calculate_file_hash(df: pd.DataFrame) -> str:
    """Calculate hash of DataFrame for session identification"""
    # Use first 100 rows + column names for hash (fast, representative)
    sample = df.head(100).to_csv(index=False)
    return hashlib.md5(sample.encode()).hexdigest()[:12]

def get_sessions_dir() -> Path:
    """Get or create sessions directory"""
    sessions_dir = Path("/mnt/user-data/outputs/sessions")
    sessions_dir.mkdir(parents=True, exist_ok=True)
    return sessions_dir

def save_session(file_hash: str, insights: Dict, filter_criteria: Dict, model_used: str):
    """Save session to file"""
    try:
        sessions_dir = get_sessions_dir()
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        session_file = sessions_dir / f"session_{file_hash}_{timestamp}.pkl"
        
        session_data = {
            'file_hash': file_hash,
            'timestamp': datetime.now().isoformat(),
            'insights': insights,
            'filter_criteria': filter_criteria,
            'model_used': model_used,
            'agent_names': list(insights.keys())
        }
        
        with open(session_file, 'wb') as f:
            pickle.dump(session_data, f)
        
        return session_file
    except Exception as e:
        print(f"Failed to save session: {e}")
        return None

def load_latest_session(file_hash: str) -> Dict:
    """Load most recent session for a given file hash"""
    try:
        sessions_dir = get_sessions_dir()
        matching_sessions = list(sessions_dir.glob(f"session_{file_hash}_*.pkl"))
        
        if not matching_sessions:
            return None
        
        # Get most recent
        latest_session = max(matching_sessions, key=lambda p: p.stat().st_mtime)
        
        with open(latest_session, 'rb') as f:
            return pickle.load(f)
    except:
        return None

def list_all_sessions() -> List[Dict]:
    """List all available sessions"""
    try:
        sessions_dir = get_sessions_dir()
        sessions = []
        
        for session_file in sessions_dir.glob("session_*.pkl"):
            try:
                with open(session_file, 'rb') as f:
                    session_data = pickle.load(f)
                    session_data['filename'] = session_file.name
                    session_data['filepath'] = str(session_file)
                    sessions.append(session_data)
            except:
                continue
        
        return sorted(sessions, key=lambda x: x['timestamp'], reverse=True)
    except:
        return []

def merge_insights(old_insights: Dict, new_insights: Dict) -> Dict:
    """Merge old and new coaching insights"""
    merged = old_insights.copy()
    merged.update(new_insights)
    return merged

def generate_analytics_context(insights: Dict, df: pd.DataFrame) -> str:
    """Generate text context for enhanced chat with DuckDB analytics"""
    
    total_agents = len(insights)
    total_calls = df['call_id'].nunique() if 'call_id' in df.columns else 0
    
    # Theme analysis
    theme_counts = {}
    priority_dist = {'high': 0, 'medium': 0, 'low': 0}
    
    for agent_data in insights.values():
        for theme in agent_data.get('coaching_themes', []):
            theme_name = theme.get('theme', '')
            theme_counts[theme_name] = theme_counts.get(theme_name, 0) + 1
            priority_dist[theme.get('priority', 'low')] += 1
    
    top_3_themes = sorted(theme_counts.items(), key=lambda x: x[1], reverse=True)[:3] if theme_counts else []
    
    # Sentiment analysis if available
    avg_sentiment = df['sentiment_score'].mean() if 'sentiment_score' in df.columns else 0
    
    context = f"""
COACHING ANALYTICS CONTEXT:

Dataset Overview:
- Total Agents Analyzed: {total_agents}
- Total Calls: {total_calls:,}
- Average Sentiment: {avg_sentiment:.2f}/5.0

Theme Distribution:
"""
    
    if top_3_themes:
        for idx, (theme, count) in enumerate(top_3_themes, 1):
            context += f"- #{idx} Theme: {theme} ({count} agents, {count/total_agents*100:.1f}%)\n"
    
    context += f"""
Priority Breakdown:
- High Priority: {priority_dist['high']} themes requiring immediate attention
- Medium Priority: {priority_dist['medium']} themes
- Low Priority: {priority_dist['low']} themes

Available DuckDB Tables:
- transcripts: call_id, agent, speaker, message, sentiment_score, timestamp
- coaching_cache: agent, theme, priority, frequency, recommendation, processed_at

You can answer questions about:
- Agent performance comparisons
- Theme correlations and patterns
- Sentiment analysis by agent or theme
- Call volume and engagement metrics
"""
    
    return context

def generate_executive_summary(insights: Dict, df: pd.DataFrame) -> str:
    """Generate executive summary section for HTML report"""
    
    total_agents = len(insights)
    total_calls = df['call_id'].nunique() if 'call_id' in df.columns else 0
    
    # Theme analysis
    theme_counts = {}
    priority_dist = {'high': 0, 'medium': 0, 'low': 0}
    
    for agent_data in insights.values():
        for theme in agent_data.get('coaching_themes', []):
            theme_name = theme.get('theme', '')
            theme_counts[theme_name] = theme_counts.get(theme_name, 0) + 1
            priority_dist[theme.get('priority', 'low')] += 1
    
    top_themes = sorted(theme_counts.items(), key=lambda x: x[1], reverse=True)[:3] if theme_counts else []
    
    # 5C scores
    c_needs = {}
    for agent_name, agent_data in insights.items():
        themes = agent_data.get('coaching_themes', [])
        for theme in themes:
            theme_name = theme.get('theme', '')
            for c, theme_list in THEME_TO_5C.items():
                if any(t.lower() in theme_name.lower() for t in theme_list):
                    c_needs[c] = c_needs.get(c, 0) + 1
    
    top_c = max(c_needs.items(), key=lambda x: x[1])[0] if c_needs else "Connection"
    
    # Sentiment
    avg_sentiment = df['sentiment_score'].mean() if 'sentiment_score' in df.columns else 0
    
    # Agent segmentation
    high_need = sum(1 for a in insights.values() if len(a.get('coaching_themes', [])) >= 4)
    moderate = sum(1 for a in insights.values() if 2 <= len(a.get('coaching_themes', [])) < 4)
    high_perf = total_agents - high_need - moderate
    
    # Generate summary HTML
    summary = f"""
<div style="background: #f8fafc; border-left: 4px solid #0ea5e9; padding: 24px; margin-bottom: 40px; border-radius: 8px;">
    <h2 style="color: #1e293b; margin-top: 0; display: flex; align-items: center; gap: 10px;">
        <svg width="24" height="24" viewBox="0 0 24 24" fill="none" stroke="#0ea5e9" stroke-width="2" stroke-linecap="round" stroke-linejoin="round">
            <path d="M14 2H6a2 2 0 0 0-2 2v16a2 2 0 0 0 2 2h12a2 2 0 0 0 2-2V8z"></path>
            <polyline points="14 2 14 8 20 8"></polyline>
            <line x1="16" y1="13" x2="8" y2="13"></line>
            <line x1="16" y1="17" x2="8" y2="17"></line>
            <polyline points="10 9 9 9 8 9"></polyline>
        </svg>
        Executive Summary
    </h2>
    
    <h3 style="color: #475569; font-size: 1.05rem; margin-top: 20px; font-weight: 600;">Analysis Overview</h3>
    <ul style="color: #64748b; line-height: 1.8; margin: 10px 0;">
        <li><strong>{total_agents}</strong> agents analyzed across <strong>{total_calls:,}</strong> calls</li>
        <li>Report generated: <strong>{datetime.now().strftime('%B %d, %Y')}</strong></li>
        <li>Average sentiment score: <strong>{avg_sentiment:.2f}/5.0</strong></li>
    </ul>
    
    <h3 style="color: #475569; font-size: 1.05rem; margin-top: 20px; font-weight: 600;">Critical Findings</h3>
    <ul style="color: #64748b; line-height: 1.8; margin: 10px 0;">
        <li><strong style="color: #dc2626;">{priority_dist['high']} agents</strong> require immediate coaching intervention (high-priority themes)</li>
"""
    
    if top_themes:
        themes_str = ", ".join([f"<strong>{theme}</strong> ({count} agents)" for theme, count in top_themes])
        summary += f"        <li>Top coaching needs: {themes_str}</li>\n"
    
    summary += f"""        <li><strong>{top_c}</strong> pillar shows highest coaching demand ({c_needs.get(top_c, 0)} agents affected)</li>
    </ul>
    
    <h3 style="color: #475569; font-size: 1.05rem; margin-top: 20px; font-weight: 600;">Agent Segmentation</h3>
    <ul style="color: #64748b; line-height: 1.8; margin: 10px 0;">
        <li><strong style="color: #10b981;">High Performers:</strong> {high_perf} agents (0-1 coaching themes)</li>
        <li><strong style="color: #f59e0b;">Moderate Development:</strong> {moderate} agents (2-3 themes)</li>
        <li><strong style="color: #dc2626;">Priority Coaching:</strong> {high_need} agents (4+ themes)</li>
    </ul>
    
    <h3 style="color: #475569; font-size: 1.05rem; margin-top: 20px; font-weight: 600;">Recommended Actions</h3>
    <ol style="color: #64748b; line-height: 1.8; margin: 10px 0;">
        <li>Immediate 1-on-1 coaching sessions for <strong>{priority_dist['high']}</strong> high-priority agents</li>
        <li>Launch targeted <strong>{top_c}</strong> training program (impacts {c_needs.get(top_c, 0)} agents)</li>
"""
    
    if top_themes:
        summary += f"        <li>Conduct focused workshop on <strong>{top_themes[0][0]}</strong> (most common challenge)</li>\n"
    
    summary += f"""        <li>Establish monthly coaching review cycle for moderate-need agents</li>
    </ol>
</div>
"""
    
    return summary

def generate_html_report(insights: Dict, df: pd.DataFrame) -> str:
    """Generate beautiful HTML report"""
    
    total_calls = len(df['call_id'].unique()) if 'call_id' in df.columns else len(df)
    total_agents = len(insights)
    total_themes = sum(len(agent_data.get('coaching_themes', [])) for agent_data in insights.values())
    avg_sentiment = df['sentiment_score'].mean() if 'sentiment_score' in df.columns else 0
    
    html = f"""
    <!DOCTYPE html>
    <html>
    <head>
        <meta charset="UTF-8">
        <meta name="viewport" content="width=device-width, initial-scale=1.0">
        <title>QA Coaching Intelligence Report</title>
        <link href="https://fonts.googleapis.com/css2?family=Inter:wght@300;400;600;700&display=swap" rel="stylesheet">
        <style>
            * {{
                margin: 0;
                padding: 0;
                box-sizing: border-box;
                font-family: 'Inter', sans-serif;
            }}
            
            body {{
                background: #ffffff;
                padding: 40px 20px;
                min-height: 100vh;
            }}
            
            .container {{
                max-width: 1200px;
                margin: 0 auto;
                background: white;
                padding: 40px;
            }}
            
            .header {{
                text-align: center;
                margin-bottom: 50px;
                padding-bottom: 30px;
                border-bottom: 1px solid #e2e8f0;
            }}
            
            h1 {{
                font-size: 2.25rem;
                font-weight: 600;
                color: #1e293b;
                margin-bottom: 10px;
            }}
            
            .subtitle {{
                color: #64748b;
                font-size: 1.1rem;
            }}
            
            .metrics-grid {{
                display: grid;
                grid-template-columns: repeat(auto-fit, minmax(200px, 1fr));
                gap: 20px;
                margin-bottom: 40px;
            }}
            
            .metric-card {{
                background: white;
                border: 1px solid #e2e8f0;
                padding: 24px;
                border-radius: 12px;
                transition: all 0.3s ease;
                box-shadow: 0 1px 3px rgba(0,0,0,0.05);
            }}
            
            .metric-card:hover {{
                transform: translateY(-2px);
                box-shadow: 0 4px 12px rgba(0,0,0,0.1);
                border-color: #0ea5e9;
            }}
            
            .metric-label {{
                font-size: 0.8rem;
                color: #64748b;
                margin-bottom: 8px;
                text-transform: uppercase;
                letter-spacing: 0.5px;
                font-weight: 600;
            }}
            
            .metric-value {{
                font-size: 2rem;
                font-weight: 700;
                color: #0ea5e9;
            }}
            
            .section-title {{
                font-size: 1.75rem;
                font-weight: 700;
                color: #1e293b;
                margin: 50px 0 30px 0;
                padding-bottom: 15px;
                border-bottom: 2px solid #e2e8f0;
                display: flex;
                align-items: center;
                gap: 12px;
            }}
            
            .chart-container {{
                background: white;
                border: 1px solid #e2e8f0;
                padding: 30px;
                border-radius: 12px;
                margin-bottom: 40px;
                box-shadow: 0 1px 3px rgba(0,0,0,0.05);
            }}
            
            .agent-grid {{
                display: grid;
                grid-template-columns: repeat(auto-fill, minmax(450px, 1fr));
                gap: 24px;
                margin-bottom: 50px;
            }}
            
            .agent-card {{
                background: white;
                border: 1px solid #e2e8f0;
                border-radius: 12px;
                padding: 30px;
                box-shadow: 0 5px 15px rgba(0,0,0,0.1);
                transition: all 0.3s ease;
            }}
            
            .agent-card:hover {{
                transform: translateY(-5px);
                box-shadow: 0 15px 35px rgba(102,126,234,0.3);
                border-color: #667eea;
            }}
            
            .agent-header {{
                display: flex;
                justify-content: space-between;
                align-items: center;
                margin-bottom: 20px;
                padding-bottom: 20px;
                border-bottom: 2px solid #f0f0f0;
            }}
            
            .agent-name {{
                font-size: 1.5rem;
                font-weight: 700;
                color: #333;
            }}
            
            .agent-stats {{
                display: flex;
                gap: 15px;
            }}
            
            .stat-badge {{
                background: #f0f0f0;
                padding: 8px 15px;
                border-radius: 10px;
                font-size: 0.85rem;
                font-weight: 600;
            }}
            
            .theme-list {{
                margin: 20px 0;
            }}
            
            .theme-item {{
                background: #f8f9fa;
                padding: 15px;
                border-radius: 10px;
                margin-bottom: 15px;
                border-left: 5px solid #667eea;
            }}
            
            .theme-item.high {{
                border-left-color: #f5576c;
                background: linear-gradient(90deg, rgba(245,87,108,0.1) 0%, transparent 100%);
            }}
            
            .theme-item.medium {{
                border-left-color: #ffa726;
                background: linear-gradient(90deg, rgba(255,167,38,0.1) 0%, transparent 100%);
            }}
            
            .theme-item.low {{
                border-left-color: #66bb6a;
                background: linear-gradient(90deg, rgba(102,187,106,0.1) 0%, transparent 100%);
            }}
            
            .theme-header {{
                display: flex;
                justify-content: space-between;
                align-items: center;
                margin-bottom: 10px;
            }}
            
            .theme-name {{
                font-weight: 700;
                font-size: 1.1rem;
                color: #333;
            }}
            
            .priority-badge {{
                padding: 5px 15px;
                border-radius: 20px;
                font-size: 0.8rem;
                font-weight: 700;
                text-transform: uppercase;
            }}
            
            .priority-high {{
                background: linear-gradient(135deg, #f093fb 0%, #f5576c 100%);
                color: white;
            }}
            
            .priority-medium {{
                background: linear-gradient(135deg, #ffd89b 0%, #ffa726 100%);
                color: white;
            }}
            
            .priority-low {{
                background: linear-gradient(135deg, #a8edea 0%, #66bb6a 100%);
                color: white;
            }}
            
            .theme-examples {{
                margin: 10px 0;
                font-size: 0.9rem;
                color: #666;
                font-style: italic;
            }}
            
            .theme-recommendation {{
                margin-top: 10px;
                padding: 10px;
                background: white;
                border-radius: 8px;
                font-size: 0.95rem;
                color: #333;
            }}
            
            .strengths-section {{
                margin-top: 20px;
                padding: 20px;
                background: linear-gradient(135deg, #11998e 0%, #38ef7d 100%);
                border-radius: 15px;
                color: white;
            }}
            
            .strengths-title {{
                font-weight: 700;
                font-size: 1.1rem;
                margin-bottom: 10px;
            }}
            
            .strengths-list {{
                list-style: none;
            }}
            
            .strengths-list li {{
                padding: 8px 0;
                padding-left: 25px;
                position: relative;
            }}
            
            .strengths-list li:before {{
                content: "✓";
                position: absolute;
                left: 0;
                font-weight: bold;
                font-size: 1.2rem;
            }}
            
            .footer {{
                text-align: center;
                margin-top: 60px;
                padding-top: 30px;
                border-top: 2px solid #e0e0e0;
                color: #666;
            }}
        </style>
    </head>
    <body>
        <div class="container">
            <div class="header">
                <h1>
                    <svg width="32" height="32" viewBox="0 0 24 24" fill="none" stroke="#0ea5e9" stroke-width="2" stroke-linecap="round" stroke-linejoin="round" style="display: inline-block; vertical-align: middle; margin-right: 10px;">
                        <circle cx="12" cy="12" r="10"></circle>
                        <path d="M12 6v6l4 2"></path>
                    </svg>
                    QA Coaching Intelligence Report
                </h1>
                <p class="subtitle">Generated on {datetime.now().strftime('%B %d, %Y at %I:%M %p')}</p>
            </div>
            
            <div class="metrics-grid">
                <div class="metric-card">
                    <div class="metric-label">Total Calls Analyzed</div>
                    <div class="metric-value">{total_calls}</div>
                </div>
                <div class="metric-card">
                    <div class="metric-label">Agents Reviewed</div>
                    <div class="metric-value">{total_agents}</div>
                </div>
                <div class="metric-card">
                    <div class="metric-label">Coaching Opportunities</div>
                    <div class="metric-value">{total_themes}</div>
                </div>
                <div class="metric-card">
                    <div class="metric-label">Average Sentiment</div>
                    <div class="metric-value">{f"{avg_sentiment:.2f}" if avg_sentiment > 0 else "N/A"}</div>
                </div>
            </div>
            
            {generate_executive_summary(insights, df)}
            
            <h2 class="section-title">
                <svg width="24" height="24" viewBox="0 0 24 24" fill="none" stroke="#0ea5e9" stroke-width="2" stroke-linecap="round" stroke-linejoin="round">
                    <path d="M21 16V8a2 2 0 0 0-1-1.73l-7-4a2 2 0 0 0-2 0l-7 4A2 2 0 0 0 3 8v8a2 2 0 0 0 1 1.73l7 4a2 2 0 0 0 2 0l7-4A2 2 0 0 0 21 16z"></path>
                    <polyline points="3.29 7 12 12 20.71 7"></polyline>
                    <line x1="12" y1="22" x2="12" y2="12"></line>
                </svg>
                5C Coaching Framework Analysis
            </h2>
            <div style="margin: 20px 0; padding: 20px; background: linear-gradient(135deg, #667eea10 0%, #764ba210 100%); border-radius: 15px;">
                <p style="font-size: 0.95rem; color: #666; text-align: center; margin-bottom: 20px;">
                    Coaching themes mapped to the 5 fundamental pillars of customer service excellence
                </p>
                <div style="display: grid; grid-template-columns: repeat(auto-fit, minmax(280px, 1fr)); gap: 20px;">
    """
    
    # Calculate 5C scores per agent
    agent_5c_scores = {}
    for agent_name, agent_data in insights.items():
        themes = agent_data.get('coaching_themes', [])
        c_scores = {c: 0 for c in FIVE_C_ICONS.keys()}
        
        for theme in themes:
            theme_name = theme.get('theme', '')
            # Find which C this theme belongs to
            for c, theme_list in THEME_TO_5C.items():
                if any(t in theme_name for t in theme_list):
                    priority = theme.get('priority', 'low')
                    # Weight by priority
                    weight = 3 if priority == 'high' else 2 if priority == 'medium' else 1
                    c_scores[c] += weight
                    break
        
        agent_5c_scores[agent_name] = c_scores
    
    # Generate 5C cards
    for c_name, icon in FIVE_C_ICONS.items():
        color = FIVE_C_COLORS[c_name]
        
        # Get top 5 agents needing help in this C (highest scores = most issues)
        agent_scores = [(agent, scores[c_name]) for agent, scores in agent_5c_scores.items()]
        top_agents = sorted(agent_scores, key=lambda x: x[1], reverse=True)[:5]
        top_agents = [a for a in top_agents if a[1] > 0]  # Only agents with issues
        
        max_score = max([score for _, score in top_agents]) if top_agents else 1
        
        html += f"""
                    <div style="background: white; border-radius: 12px; padding: 18px; box-shadow: 0 4px 12px rgba(0,0,0,0.08); transition: all 0.3s ease;">
                        <div style="display: flex; align-items: center; gap: 12px; margin-bottom: 15px; padding-bottom: 12px; border-bottom: 2px solid {color};">
                            <div>{icon}</div>
                            <div>
                                <h3 style="font-size: 1.1rem; font-weight: 700; color: #333; margin: 0;">{c_name}</h3>
                                <p style="font-size: 0.8rem; color: #666; margin: 3px 0 0 0;">{len([a for a, s in top_agents])} agents need support</p>
                            </div>
                        </div>
        """
        
        if top_agents:
            html += """
                        <div style="display: flex; flex-direction: column; gap: 10px;">
            """
            
            for idx, (agent, score) in enumerate(top_agents, 1):
                percentage = (score / max_score) * 100
                
                # Get specific themes for this agent in this C
                agent_themes = insights[agent].get('coaching_themes', [])
                relevant_themes = []
                for theme in agent_themes:
                    theme_name = theme.get('theme', '')
                    if any(t in theme_name for t in THEME_TO_5C[c_name]):
                        relevant_themes.append(theme_name)
                
                themes_text = ', '.join(relevant_themes[:2])
                if len(relevant_themes) > 2:
                    themes_text += f" +{len(relevant_themes) - 2}"
                
                html += f"""
                            <div style="background: #f8f9fa; padding: 10px; border-radius: 8px;">
                                <div style="display: flex; justify-content: space-between; align-items: center; margin-bottom: 6px;">
                                    <div style="display: flex; align-items: center; gap: 8px;">
                                        <span style="font-weight: 700; color: #333; font-size: 0.9rem;">{idx}.</span>
                                        <span style="font-weight: 600; color: #333; font-size: 0.9rem;">{agent}</span>
                                    </div>
                                    <span style="font-size: 0.75rem; color: #666; font-weight: 600;">Score: {score}</span>
                                </div>
                                <div style="background: #e0e0e0; height: 6px; border-radius: 8px; overflow: hidden; margin-bottom: 6px;">
                                    <div style="background: {color}; height: 100%; width: {percentage}%; border-radius: 8px; transition: width 0.5s ease;"></div>
                                </div>
                                <div style="font-size: 0.75rem; color: #666; font-style: italic;">{themes_text}</div>
                            </div>
                """
            
            html += """
                        </div>
            """
        else:
            html += """
                        <div style="text-align: center; padding: 20px; color: #999;">
                            <p style="font-size: 0.95rem; margin: 0;">✨ Great job!</p>
                            <p style="font-size: 0.8rem; margin: 5px 0 0 0;">No major issues</p>
                        </div>
            """
        
        html += """
                    </div>
        """
    
    html += """
                </div>
            </div>
            
            <h2 class="section-title">
                <svg width="24" height="24" viewBox="0 0 24 24" fill="none" stroke="#0ea5e9" stroke-width="2" stroke-linecap="round" stroke-linejoin="round">
                    <rect x="3" y="3" width="18" height="18" rx="2" ry="2"></rect>
                    <line x1="3" y1="9" x2="21" y2="9"></line>
                    <line x1="9" y1="21" x2="9" y2="9"></line>
                </svg>
                Coaching Theme Distribution
            </h2>
            <div class="chart-container">
                <div style="display: grid; grid-template-columns: repeat(auto-fill, minmax(280px, 1fr)); gap: 25px; padding: 20px;">
    """
    
    # Prepare theme data with icons
    theme_icons = {
        "Active Listening": '<svg width="20" height="20" viewBox="0 0 24 24" fill="none" stroke="#64748b" stroke-width="2"><circle cx="12" cy="12" r="1"/><path d="M12 1v6m0 6v6m-6-6h6m6 0h6"/></svg>',
        "Empathy": '<svg width="20" height="20" viewBox="0 0 24 24" fill="none" stroke="#64748b" stroke-width="2"><path d="M20.84 4.61a5.5 5.5 0 0 0-7.78 0L12 5.67l-1.06-1.06a5.5 5.5 0 0 0-7.78 7.78l1.06 1.06L12 21.23l7.78-7.78 1.06-1.06a5.5 5.5 0 0 0 0-7.78z"/></svg>',
        "Communication": '<svg width="20" height="20" viewBox="0 0 24 24" fill="none" stroke="#64748b" stroke-width="2"><path d="M21 15a2 2 0 0 1-2 2H7l-4 4V5a2 2 0 0 1 2-2h14a2 2 0 0 1 2 2z"/></svg>',
        "Professional": '<svg width="20" height="20" viewBox="0 0 24 24" fill="none" stroke="#64748b" stroke-width="2"><rect x="2" y="7" width="20" height="14" rx="2"/><path d="M16 21V5a2 2 0 0 0-2-2h-4a2 2 0 0 0-2 2v16"/></svg>',
        "Resolution": '<svg width="20" height="20" viewBox="0 0 24 24" fill="none" stroke="#64748b" stroke-width="2"><polyline points="20 6 9 17 4 12"/></svg>',
        "Problem": '<svg width="20" height="20" viewBox="0 0 24 24" fill="none" stroke="#64748b" stroke-width="2"><circle cx="11" cy="11" r="8"/><path d="m21 21-4.35-4.35"/></svg>',
        "Solution": '<svg width="20" height="20" viewBox="0 0 24 24" fill="none" stroke="#64748b" stroke-width="2"><path d="M15 3h4a2 2 0 0 1 2 2v14a2 2 0 0 1-2 2h-4m-5-4 3 3 3-3m-3-10v12"/></svg>',
        "Response Time": '<svg width="20" height="20" viewBox="0 0 24 24" fill="none" stroke="#64748b" stroke-width="2"><circle cx="12" cy="12" r="10"/><polyline points="12 6 12 12 16 14"/></svg>',
        "Process": '<svg width="20" height="20" viewBox="0 0 24 24" fill="none" stroke="#64748b" stroke-width="2"><path d="M14 2H6a2 2 0 0 0-2 2v16a2 2 0 0 0 2 2h12a2 2 0 0 0 2-2V8z"/><polyline points="14 2 14 8 20 8"/></svg>',
        "Escalation": '<svg width="20" height="20" viewBox="0 0 24 24" fill="none" stroke="#64748b" stroke-width="2"><polyline points="18 15 12 9 6 15"/></svg>',
        "Proactive": '<svg width="20" height="20" viewBox="0 0 24 24" fill="none" stroke="#64748b" stroke-width="2"><circle cx="12" cy="12" r="10"/><circle cx="12" cy="12" r="3"/></svg>',
        "Expectations": '<svg width="20" height="20" viewBox="0 0 24 24" fill="none" stroke="#64748b" stroke-width="2"><line x1="18" y1="20" x2="18" y2="10"/><line x1="12" y1="20" x2="12" y2="4"/><line x1="6" y1="20" x2="6" y2="14"/></svg>',
        "Difficult": '<svg width="20" height="20" viewBox="0 0 24 24" fill="none" stroke="#64748b" stroke-width="2"><path d="M10.29 3.86L1.82 18a2 2 0 0 0 1.71 3h16.94a2 2 0 0 0 1.71-3L13.71 3.86a2 2 0 0 0-3.42 0z"/><line x1="12" y1="9" x2="12" y2="13"/><line x1="12" y1="17" x2="12.01" y2="17"/></svg>',
        "Rapport": '<svg width="20" height="20" viewBox="0 0 24 24" fill="none" stroke="#64748b" stroke-width="2"><path d="M17 21v-2a4 4 0 0 0-4-4H5a4 4 0 0 0-4 4v2"/><circle cx="9" cy="7" r="4"/><path d="M23 21v-2a4 4 0 0 0-3-3.87"/></svg>',
        "Knowledge": '<svg width="20" height="20" viewBox="0 0 24 24" fill="none" stroke="#64748b" stroke-width="2"><path d="M4 19.5A2.5 2.5 0 0 1 6.5 17H20"/><path d="M6.5 2H20v20H6.5A2.5 2.5 0 0 1 4 19.5v-15A2.5 2.5 0 0 1 6.5 2z"/></svg>',
        "Confidence": '<svg width="20" height="20" viewBox="0 0 24 24" fill="none" stroke="#64748b" stroke-width="2"><polyline points="22 12 18 12 15 21 9 3 6 12 2 12"/></svg>'
    }
    
    theme_counts = {}
    for agent_data in insights.values():
        for theme in agent_data.get('coaching_themes', []):
            theme_name = theme.get('theme', '')
            theme_counts[theme_name] = theme_counts.get(theme_name, 0) + 1
    
    sorted_themes = sorted(theme_counts.items(), key=lambda x: x[1], reverse=True)[:8]
    max_count = max([c for _, c in sorted_themes]) if sorted_themes else 1
    
    for theme_name, count in sorted_themes:
        # Find matching icon
        icon = '<svg width="20" height="20" viewBox="0 0 24 24" fill="none" stroke="#64748b" stroke-width="2"><circle cx="12" cy="12" r="10"/><circle cx="12" cy="12" r="3"/></svg>'
        for key, svg_icon in theme_icons.items():
            if key.lower() in theme_name.lower():
                icon = svg_icon
                break
        
        percentage = (count / max_count) * 100
        
        html += f"""
                    <div style="background: white; border: 1px solid #e2e8f0; border-radius: 12px; padding: 20px; transition: all 0.3s ease;">
                        <div style="display: flex; align-items: center; gap: 12px; margin-bottom: 12px;">
                            <div>{icon}</div>
                            <div style="flex: 1;">
                                <div style="font-weight: 600; font-size: 1rem; color: #1e293b; margin-bottom: 4px;">{theme_name}</div>
                                <div style="font-size: 0.85rem; color: #64748b;">{count} agents</div>
                            </div>
                        </div>
                        <div style="background: #f1f5f9; height: 8px; border-radius: 6px; overflow: hidden;">
                            <div style="background: #0ea5e9; height: 100%; width: {percentage}%; border-radius: 6px; transition: width 0.5s ease;"></div>
                        </div>
                    </div>
        """
    
    html += """
                </div>
            </div>
            
            <h2 class="section-title">
                <svg width="24" height="24" viewBox="0 0 24 24" fill="none" stroke="#0ea5e9" stroke-width="2" stroke-linecap="round" stroke-linejoin="round">
                    <line x1="12" y1="1" x2="12" y2="23"></line>
                    <path d="M17 5H9.5a3.5 3.5 0 0 0 0 7h5a3.5 3.5 0 0 1 0 7H6"></path>
                </svg>
                Agent Performance Summary
            </h2>
            <div style="margin: 20px 0;">
    """
    
    # Sort agents by priority (high priority themes first)
    agent_priority = []
    for agent_name, agent_data in insights.items():
        themes = agent_data.get('coaching_themes', [])
        high_priority = sum(1 for t in themes if t.get('priority') == 'high')
        medium_priority = sum(1 for t in themes if t.get('priority') == 'medium')
        agent_priority.append((agent_name, high_priority, medium_priority, agent_data))
    
    agent_priority.sort(key=lambda x: (x[1], x[2]), reverse=True)  # Sort by high, then medium
    
    total_agents = len(agent_priority)
    agents_per_page = 20
    total_pages = (total_agents + agents_per_page - 1) // agents_per_page
    
    html += f"""
                <div style="background: #f8f9fa; padding: 15px; border-radius: 10px; margin-bottom: 20px;">
                    <p style="margin: 0; color: #666; font-size: 0.95rem;">
                        <strong>Showing {total_agents} agents</strong> | Sorted by priority | {agents_per_page} per page
                    </p>
                </div>
    """
    
    # Generate pagination controls
    for page in range(total_pages):
        page_num = page + 1
        start_idx = page * agents_per_page
        end_idx = min(start_idx + agents_per_page, total_agents)
        
        page_agents = agent_priority[start_idx:end_idx]
        
        html += f"""
                <div id="page-{page_num}" style="display: {'block' if page == 0 else 'none'};">
                    <div style="overflow-x: auto;">
                        <table style="width: 100%; border-collapse: separate; border-spacing: 0 12px;">
                            <thead>
                                <tr style="background: #f8fafc; border-bottom: 2px solid #e2e8f0;">
                                    <th style="padding: 15px; text-align: left; color: #475569; font-weight: 600; font-size: 0.85rem; border-radius: 8px 0 0 0; text-transform: uppercase; letter-spacing: 0.5px;">Agent</th>
                                    <th style="padding: 15px; text-align: left; color: #475569; font-weight: 600; font-size: 0.85rem; text-transform: uppercase; letter-spacing: 0.5px;">Calls</th>
                                    <th style="padding: 15px; text-align: left; color: #475569; font-weight: 600; font-size: 0.85rem; text-transform: uppercase; letter-spacing: 0.5px;">Top Improvement Area</th>
                                    <th style="padding: 15px; text-align: left; color: #475569; font-weight: 600; font-size: 0.85rem; border-radius: 0 8px 0 0; text-transform: uppercase; letter-spacing: 0.5px;">Priority</th>
                                </tr>
                            </thead>
                            <tbody>
        """
        
        for agent_name, high_pri, med_pri, agent_data in page_agents:
            themes = agent_data.get('coaching_themes', [])
            calls = agent_data.get('calls_analyzed', 0)
            
            if themes:
                top_theme = themes[0]
                theme_name = top_theme.get('theme', 'N/A')
                priority = top_theme.get('priority', 'low')
                
                # Get icon for theme
                icon = '<svg width="16" height="16" viewBox="0 0 24 24" fill="none" stroke="#64748b" stroke-width="2"><circle cx="12" cy="12" r="10"/><circle cx="12" cy="12" r="3"/></svg>'
                for key, svg_icon in theme_icons.items():
                    if key.lower() in theme_name.lower():
                        icon = svg_icon
                        break
                
                # Priority colors
                if priority == 'high':
                    priority_color = "#dc2626"  # Red
                    priority_icon = ""
                elif priority == 'medium':
                    priority_color = "#f59e0b"  # Amber
                    priority_icon = ""
                else:
                    priority_color = "#10b981"  # Green
                    priority_icon = ""
                
                html += f"""
                                <tr style="background: white; box-shadow: 0 2px 8px rgba(0,0,0,0.04); transition: all 0.3s ease;">
                                    <td style="padding: 15px; font-weight: 700; font-size: 0.95rem; color: #333; border-radius: 10px 0 0 10px;">
                                        <div style="display: flex; align-items: center; gap: 10px;">
                                            <div style="width: 38px; height: 38px; border-radius: 50%; background: #0ea5e9; display: flex; align-items: center; justify-content: center; color: white; font-weight: 700; font-size: 1rem;">
                                                {agent_name[0].upper()}
                                            </div>
                                            {agent_name}
                                        </div>
                                    </td>
                                    <td style="padding: 15px;">
                                        <div style="display: inline-block; background: #f1f5f9; padding: 6px 14px; border-radius: 15px; font-weight: 600; color: #475569; font-size: 0.85rem;">
                                            {calls} calls
                                        </div>
                                    </td>
                                    <td style="padding: 15px;">
                                        <div style="display: inline-flex; align-items: center; gap: 8px; background: #f8fafc; padding: 8px 16px; border-radius: 16px; border: 1px solid #e2e8f0;">
                                            <span>{icon}</span>
                                            <span style="font-weight: 600; color: #1e293b; font-size: 0.9rem;">{theme_name}</span>
                                        </div>
                                    </td>
                                    <td style="padding: 15px; border-radius: 0 10px 10px 0;">
                                        <div style="display: inline-flex; align-items: center; gap: 6px; background: {priority_color}; padding: 8px 16px; border-radius: 20px; color: white; font-weight: 700; text-transform: uppercase; font-size: 0.8rem;">
                                            <span>{priority_icon}</span>
                                            <span>{priority}</span>
                                        </div>
                                    </td>
                                </tr>
                """
            else:
                html += f"""
                                <tr style="background: white; box-shadow: 0 2px 8px rgba(0,0,0,0.04);">
                                    <td style="padding: 15px; font-weight: 700; font-size: 0.95rem; color: #333; border-radius: 10px 0 0 10px;">
                                        <div style="display: flex; align-items: center; gap: 10px;">
                                            <div style="width: 38px; height: 38px; border-radius: 50%; background: #0ea5e9; display: flex; align-items: center; justify-content: center; color: white; font-weight: 700; font-size: 1rem;">
                                                {agent_name[0].upper()}
                                            </div>
                                            {agent_name}
                                        </div>
                                    </td>
                                    <td style="padding: 15px;">
                                        <div style="display: inline-block; background: #f0f0f0; padding: 6px 14px; border-radius: 15px; font-weight: 600; color: #666; font-size: 0.85rem;">
                                            {calls} calls
                                        </div>
                                    </td>
                                    <td style="padding: 15px;" colspan="2">
                                        <div style="color: #999; font-style: italic; font-size: 0.9rem;">No coaching themes identified</div>
                                    </td>
                                </tr>
                """
        
        html += """
                            </tbody>
                        </table>
                    </div>
                </div>
        """
    
    # Add pagination controls
    if total_pages > 1:
        html += """
                <div style="display: flex; justify-content: center; align-items: center; gap: 10px; margin: 30px 0;">
        """
        
        for page in range(total_pages):
            page_num = page + 1
            html += f"""
                    <button onclick="showPage({page_num})" id="page-btn-{page_num}" style="padding: 10px 16px; border: 2px solid #0ea5e9; background: {'#0ea5e9' if page == 0 else 'white'}; color: {'white' if page == 0 else '#0ea5e9'}; border-radius: 8px; font-weight: 600; cursor: pointer; transition: all 0.3s ease; font-size: 0.9rem;">
                        {page_num}
                    </button>
            """
        
        html += """
                </div>
                
                <script>
                function showPage(pageNum) {
                    // Hide all pages
                    const pages = document.querySelectorAll('[id^="page-"]');
                    pages.forEach(page => {
                        if (page.id.startsWith('page-btn-')) return;
                        page.style.display = 'none';
                    });
                    
                    // Show selected page
                    document.getElementById('page-' + pageNum).style.display = 'block';
                    
                    // Update button styles
                    const buttons = document.querySelectorAll('[id^="page-btn-"]');
                    buttons.forEach(btn => {
                        btn.style.background = 'white';
                        btn.style.color = '#0ea5e9';
                    });
                    
                    document.getElementById('page-btn-' + pageNum).style.background = '#0ea5e9';
                    document.getElementById('page-btn-' + pageNum).style.color = 'white';
                }
                </script>
        """
    
    html += """
            </div>
            
            <h2 class="section-title">
                <svg width="24" height="24" viewBox="0 0 24 24" fill="none" stroke="#0ea5e9" stroke-width="2" stroke-linecap="round" stroke-linejoin="round">
                    <path d="M17 21v-2a4 4 0 0 0-4-4H5a4 4 0 0 0-4 4v2"></path>
                    <circle cx="9" cy="7" r="4"></circle>
                    <path d="M23 21v-2a4 4 0 0 0-3-3.87"></path>
                    <path d="M16 3.13a4 4 0 0 1 0 7.75"></path>
                </svg>
                Agent Coaching Details
            </h2>
            <div class="agent-grid">
    """
    
    # Add agent cards
    for agent_name, agent_data in insights.items():
        themes = agent_data.get('coaching_themes', [])
        strengths = agent_data.get('strengths', [])
        calls_analyzed = agent_data.get('calls_analyzed', 0)
        
        # Generate email share link
        mailto_link = generate_email_share_link(agent_name, agent_data)
        
        html += f"""
                <div class="agent-card">
                    <div class="agent-header" style="display: flex; justify-content: space-between; align-items: flex-start; margin-bottom: 20px;">
                        <div style="flex: 1;">
                            <div class="agent-name">
                                <svg width="20" height="20" viewBox="0 0 24 24" fill="none" stroke="#64748b" stroke-width="2" stroke-linecap="round" stroke-linejoin="round" style="display: inline-block; vertical-align: middle; margin-right: 8px;">
                                    <path d="M20 21v-2a4 4 0 0 0-4-4H8a4 4 0 0 0-4 4v2"></path>
                                    <circle cx="12" cy="7" r="4"></circle>
                                </svg>
                                {agent_name}
                            </div>
                            <div class="agent-stats">
                                <div class="stat-badge">{calls_analyzed} calls</div>
                                <div class="stat-badge">{len(themes)} themes</div>
                            </div>
                        </div>
                        <a href="{mailto_link}" title="Share via email" style="flex-shrink: 0; width: 36px; height: 36px; background: #ffffff; border: 1px solid #e2e8f0; border-radius: 50%; display: flex; align-items: center; justify-content: center; text-decoration: none; transition: all 0.3s ease; margin-left: 16px;" 
                           onmouseover="this.style.borderColor='#0ea5e9'; this.style.boxShadow='0 2px 6px rgba(14, 165, 233, 0.3)';" 
                           onmouseout="this.style.borderColor='#e2e8f0'; this.style.boxShadow='none';">
                            <svg width="16" height="16" viewBox="0 0 24 24" fill="none" stroke="#64748b" stroke-width="2" stroke-linecap="round" stroke-linejoin="round">
                                <path d="M4 4h16c1.1 0 2 .9 2 2v12c0 1.1-.9 2-2 2H4c-1.1 0-2-.9-2-2V6c0-1.1.9-2 2-2z"></path>
                                <polyline points="22,6 12,13 2,6"></polyline>
                            </svg>
                        </a>
                    </div>
                    
                    <div class="theme-list">
        """
        
        for theme in themes:
            priority = theme.get('priority', 'low')
            html += f"""
                        <div class="theme-item {priority}">
                            <div class="theme-header">
                                <div class="theme-name">{theme.get('theme', '')}</div>
                                <div class="priority-badge priority-{priority}">{priority}</div>
                            </div>
                            <div class="theme-examples">
                                Frequency: {theme.get('frequency', 0)} instances<br>
                                Examples: {' | '.join(theme.get('examples', [])[:2])}
                            </div>
                            <div class="theme-recommendation">
                                <strong>Recommendation:</strong> {theme.get('recommendation', '')}
                            </div>
                        </div>
            """
        
        if strengths:
            html += f"""
                    </div>
                    <div class="strengths-section">
                        <div class="strengths-title">Strengths</div>
                        <ul class="strengths-list">
            """
            for strength in strengths:
                html += f"<li>{strength}</li>"
            html += """
                        </ul>
                    </div>
            """
        else:
            html += "</div>"
        
        html += "</div>"
    
    # Close HTML
    html += """
            </div>
            
            <div class="footer">
                <p>QA Coaching Intelligence Platform | Developed by CE INNOVATIONS LAB 2025</p>
            </div>
        </div>
    </body>
    </html>
    """
    
    return html

def generate_powerpoint(insights: Dict, df: pd.DataFrame) -> bytes:
    """Generate PowerPoint presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Add gradient background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(102, 126, 234)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "QA Coaching Intelligence Report"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = f"Generated on {datetime.now().strftime('%B %d, %Y')}"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(20)
    subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Key findings slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title only
    title = slide.shapes.title
    title.text = "Key Findings"
    
    total_calls = len(df['call_id'].unique()) if 'call_id' in df.columns else len(df)
    total_agents = len(insights)
    total_themes = sum(len(agent_data.get('coaching_themes', [])) for agent_data in insights.values())
    
    # Add metrics
    metrics_text = f"""
    ✓ Analyzed {total_calls} calls
    ✓ Reviewed {total_agents} agents
    ✓ Identified {total_themes} coaching opportunities
    """
    
    text_box = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(6), Inches(3))
    text_frame = text_box.text_frame
    text_frame.text = metrics_text
    for para in text_frame.paragraphs:
        para.font.size = Pt(28)
        para.space_before = Pt(20)
    
    # 5C Framework Overview Slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "🎯 5C Coaching Framework"
    
    # Calculate 5C scores per agent
    agent_5c_scores = {}
    for agent_name, agent_data in insights.items():
        themes = agent_data.get('coaching_themes', [])
        c_scores = {c: 0 for c in FIVE_C_ICONS.keys()}
        
        for theme in themes:
            theme_name = theme.get('theme', '')
            for c, theme_list in THEME_TO_5C.items():
                if any(t in theme_name for t in theme_list):
                    priority = theme.get('priority', 'low')
                    weight = 3 if priority == 'high' else 2 if priority == 'medium' else 1
                    c_scores[c] += weight
                    break
        
        agent_5c_scores[agent_name] = c_scores
    
    # Add 5C summary
    y_pos = 2
    for c_name, icon in FIVE_C_ICONS.items():
        agent_scores = [(agent, scores[c_name]) for agent, scores in agent_5c_scores.items()]
        top_agents = sorted(agent_scores, key=lambda x: x[1], reverse=True)[:3]
        agents_needing_help = len([a for a, s in agent_scores if s > 0])
        
        text_box = slide.shapes.add_textbox(Inches(1.5), Inches(y_pos), Inches(7), Inches(0.8))
        text_frame = text_box.text_frame
        text_frame.text = f"{icon} {c_name}: {agents_needing_help} agents need support"
        para = text_frame.paragraphs[0]
        para.font.size = Pt(20)
        para.font.bold = True
        
        y_pos += 0.9
    
    # Agent slides (limit to 10)
    for agent_name, agent_data in list(insights.items())[:10]:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = f"👤 {agent_name}"
        
        # Agent's 5C scores
        agent_scores = agent_5c_scores.get(agent_name, {})
        
        # Add 5C breakdown
        y_pos = 2
        text_box = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(4), Inches(0.5))
        text_frame = text_box.text_frame
        text_frame.text = "5C Framework Scores:"
        text_frame.paragraphs[0].font.size = Pt(18)
        text_frame.paragraphs[0].font.bold = True
        
        y_pos += 0.6
        for c_name, score in agent_scores.items():
            if score > 0:
                icon = FIVE_C_ICONS[c_name]
                text_box = slide.shapes.add_textbox(Inches(1.2), Inches(y_pos), Inches(3.5), Inches(0.4))
                text_frame = text_box.text_frame
                text_frame.text = f"{icon} {c_name}: {score} issues"
                text_frame.paragraphs[0].font.size = Pt(16)
                y_pos += 0.5
        
        # Top themes
        themes = agent_data.get('coaching_themes', [])[:3]
        
        y_pos = 2
        theme_box = slide.shapes.add_textbox(Inches(5.5), Inches(y_pos), Inches(4), Inches(0.5))
        theme_frame = theme_box.text_frame
        theme_frame.text = "Top Coaching Needs:"
        theme_frame.paragraphs[0].font.size = Pt(18)
        theme_frame.paragraphs[0].font.bold = True
        
        y_pos += 0.6
        for idx, theme in enumerate(themes, 1):
            theme_box = slide.shapes.add_textbox(Inches(5.5), Inches(y_pos), Inches(4), Inches(1))
            theme_frame = theme_box.text_frame
            
            theme_frame.text = f"{idx}. {theme.get('theme', '')}"
            theme_para = theme_frame.paragraphs[0]
            theme_para.font.size = Pt(16)
            theme_para.font.bold = True
            
            priority = theme.get('priority', 'low')
            priority_text = theme_frame.add_paragraph()
            priority_text.text = f"Priority: {priority.upper()}"
            priority_text.font.size = Pt(14)
            
            if priority == 'high':
                priority_text.font.color.rgb = RGBColor(245, 87, 108)
            elif priority == 'medium':
                priority_text.font.color.rgb = RGBColor(255, 167, 38)
            else:
                priority_text.font.color.rgb = RGBColor(102, 187, 106)
            
            rec_text = theme_frame.add_paragraph()
            rec_text.text = f"💡 {theme.get('recommendation', '')[:80]}..."
            rec_text.font.size = Pt(12)
            rec_text.font.italic = True
            
            y_pos += 1.3
    
    # Next steps slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "📋 Next Steps"
    
    next_steps = """
    1. Review individual agent 5C scores
    2. Schedule coaching sessions for high-priority themes
    3. Focus on Connection & Clarity first (foundational)
    4. Track improvement over next 30 days
    5. Re-analyze to measure progress
    """
    
    text_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(6), Inches(3))
    text_frame = text_box.text_frame
    text_frame.text = next_steps
    for para in text_frame.paragraphs:
        para.font.size = Pt(24)
        para.space_before = Pt(15)
    
    # Save to bytes
    ppt_bytes = io.BytesIO()
    prs.save(ppt_bytes)
    ppt_bytes.seek(0)
    return ppt_bytes.getvalue()


# Sidebar
with st.sidebar:
    st.markdown("""
        <div style='display: flex; align-items: center; gap: 10px; margin-bottom: 20px;'>
            <svg width="24" height="24" viewBox="0 0 24 24" fill="none" stroke="#0ea5e9" stroke-width="2" stroke-linecap="round" stroke-linejoin="round">
                <path d="M12 20h9"></path>
                <path d="M16.5 3.5a2.121 2.121 0 0 1 3 3L7 19l-4 1 1-4L16.5 3.5z"></path>
            </svg>
            <h3 style='margin: 0; color: #1e293b;'>QA Coaching Intelligence</h3>
        </div>
    """, unsafe_allow_html=True)
    st.markdown("---")
    
    st.markdown("""
        <div style='display: flex; align-items: center; gap: 8px; margin-bottom: 10px;'>
            <svg width="20" height="20" viewBox="0 0 24 24" fill="none" stroke="#64748b" stroke-width="2" stroke-linecap="round" stroke-linejoin="round">
                <rect x="2" y="3" width="20" height="14" rx="2" ry="2"></rect>
                <line x1="8" y1="21" x2="16" y2="21"></line>
                <line x1="12" y1="17" x2="12" y2="21"></line>
            </svg>
            <h4 style='margin: 0; color: #475569;'>LLM Provider</h4>
        </div>
    """, unsafe_allow_html=True)
    llm_provider = st.radio(
        "Select provider:",
        ["OpenRouter", "Local LLM (LM Studio/Ollama)"],
        key="llm_provider_radio"
    )
    
    if llm_provider == "OpenRouter":
        st.session_state.llm_provider = "openrouter"
        
        # Try to get API key from secrets first
        try:
            api_key = st.secrets.get("OPENROUTER_API_KEY", "")
            if api_key:
                st.session_state.openrouter_api_key = api_key
                st.success("✅ API key loaded from secrets")
            else:
                # Fallback to manual input
                api_key = st.text_input("OpenRouter API Key:", type="password", key="api_key_input")
                st.session_state.openrouter_api_key = api_key
                if not api_key:
                    st.warning("⚠️ API key required")
        except:
            # No secrets available, use manual input
            api_key = st.text_input("OpenRouter API Key:", type="password", key="api_key_input")
            st.session_state.openrouter_api_key = api_key
            if not api_key:
                st.warning("⚠️ API key required")
    else:
        st.session_state.llm_provider = "local"
        local_url = st.text_input(
            "Local LLM URL:",
            value="http://localhost:1234/v1/chat/completions",
            key="local_llm_url_input"
        )
        st.session_state.local_llm_url = local_url
        st.info("💡 Make sure LM Studio or Ollama is running")
    
    st.markdown("---")
    st.markdown("""
        <div style='display: flex; align-items: center; gap: 8px; margin-bottom: 10px;'>
            <svg width="20" height="20" viewBox="0 0 24 24" fill="none" stroke="#64748b" stroke-width="2" stroke-linecap="round" stroke-linejoin="round">
                <circle cx="12" cy="12" r="10"></circle>
                <circle cx="12" cy="12" r="3"></circle>
                <line x1="12" y1="1" x2="12" y2="3"></line>
                <line x1="12" y1="21" x2="12" y2="23"></line>
                <line x1="4.22" y1="4.22" x2="5.64" y2="5.64"></line>
                <line x1="18.36" y1="18.36" x2="19.78" y2="19.78"></line>
                <line x1="1" y1="12" x2="3" y2="12"></line>
                <line x1="21" y1="12" x2="23" y2="12"></line>
                <line x1="4.22" y1="19.78" x2="5.64" y2="18.36"></line>
                <line x1="18.36" y1="5.64" x2="19.78" y2="4.22"></line>
            </svg>
            <h4 style='margin: 0; color: #475569;'>Analysis Model</h4>
        </div>
    """, unsafe_allow_html=True)
    
    if st.session_state.llm_provider == "openrouter":
        analysis_model = st.selectbox(
            "Select model:",
            options=list(MODELS.keys()),
            format_func=lambda x: f"[{MODELS[x]['rating']}] {MODELS[x]['name']}" + (" (Recommended)" if MODELS[x].get('recommended') else ""),
            index=0
        )
        st.info(f"**Best for:** {MODELS[analysis_model]['best_for']}\n\n**Speed:** {MODELS[analysis_model]['speed']}")
    else:
        analysis_model = st.text_input(
            "Model name:",
            value="llama3",
            help="Enter the model name from LM Studio/Ollama"
        )
    
    st.markdown("---")
    st.markdown("""
        <div style='display: flex; align-items: center; gap: 8px; margin-bottom: 10px;'>
            <svg width="20" height="20" viewBox="0 0 24 24" fill="none" stroke="#64748b" stroke-width="2" stroke-linecap="round" stroke-linejoin="round">
                <circle cx="12" cy="12" r="3"></circle>
                <path d="M12 1v6m0 6v6m5.656-14.656l-4.242 4.242m-2.828 2.828l-4.242 4.242M23 12h-6m-6 0H1m17.656 5.656l-4.242-4.242m-2.828-2.828l-4.242-4.242"></path>
            </svg>
            <h4 style='margin: 0; color: #475569;'>Processing Settings</h4>
        </div>
    """, unsafe_allow_html=True)
    
    max_concurrent = st.slider(
        "Concurrent requests:",
        min_value=5,
        max_value=20,
        value=10,
        help="More = faster but may hit rate limits"
    )
    st.session_state.max_concurrent = max_concurrent
    
    calls_per_minute = st.number_input(
        "Rate limit (calls/min):",
        min_value=10,
        max_value=100,
        value=50,
        help="Your OpenRouter/API rate limit"
    )
    st.session_state.calls_per_minute = calls_per_minute
    
    # Session Management
    st.markdown("---")
    st.markdown("""
        <div style='display: flex; align-items: center; gap: 8px; margin-bottom: 10px;'>
            <svg width="20" height="20" viewBox="0 0 24 24" fill="none" stroke="#64748b" stroke-width="2" stroke-linecap="round" stroke-linejoin="round">
                <path d="M21 15v4a2 2 0 0 1-2 2H5a2 2 0 0 1-2-2v-4"></path>
                <polyline points="7 10 12 15 17 10"></polyline>
                <line x1="12" y1="15" x2="12" y2="3"></line>
            </svg>
            <h4 style='margin: 0; color: #475569;'>Session Management</h4>
        </div>
    """, unsafe_allow_html=True)
    
    # Check for existing session if data is loaded
    if st.session_state.get('processed_df') is not None:
        file_hash = calculate_file_hash(st.session_state.processed_df)
        st.session_state.file_hash = file_hash
        
        existing_session = load_latest_session(file_hash)
        
        if existing_session:
            st.info(f"**Session Found**\n\n"
                   f"Agents: {len(existing_session['agent_names'])}\n\n"
                   f"Date: {datetime.fromisoformat(existing_session['timestamp']).strftime('%b %d, %I:%M %p')}\n\n"
                   f"Model: {existing_session['model_used'][:30]}")
            
            col1, col2 = st.columns(2)
            with col1:
                if st.button("Resume", use_container_width=True):
                    st.session_state.coaching_insights = existing_session['insights']
                    st.session_state.processed = True
                    st.session_state.analytics_context = generate_analytics_context(
                        existing_session['insights'],
                        st.session_state.processed_df
                    )
                    st.rerun()
            with col2:
                if st.button("New", use_container_width=True):
                    st.session_state.coaching_insights = {}
                    st.session_state.processed = False
                    st.rerun()
    
    # Session Browser
    with st.expander("Previous Sessions"):
        all_sessions = list_all_sessions()
        
        if all_sessions:
            for session in all_sessions[:5]:  # Show last 5
                st.write(f"**{datetime.fromisoformat(session['timestamp']).strftime('%b %d %I:%M %p')}**")
                st.write(f"{len(session['agent_names'])} agents")
                
                col1, col2 = st.columns([3, 1])
                with col1:
                    if st.button("Load", key=f"load_{session['filename']}", use_container_width=True):
                        st.session_state.coaching_insights = session['insights']
                        st.session_state.processed = True
                        st.rerun()
                with col2:
                    if st.button("🗑", key=f"del_{session['filename']}"):
                        Path(session['filepath']).unlink()
                        st.rerun()
                st.divider()
        else:
            st.caption("No saved sessions")
    
    st.markdown("---")
    st.markdown("""
        <div style='display: flex; align-items: center; gap: 8px; margin-bottom: 10px;'>
            <svg width="20" height="20" viewBox="0 0 24 24" fill="none" stroke="#64748b" stroke-width="2" stroke-linecap="round" stroke-linejoin="round">
                <path d="M14 2H6a2 2 0 0 0-2 2v16a2 2 0 0 0 2 2h12a2 2 0 0 0 2-2V8z"></path>
                <polyline points="14 2 14 8 20 8"></polyline>
                <line x1="16" y1="13" x2="8" y2="13"></line>
                <line x1="16" y1="17" x2="8" y2="17"></line>
                <polyline points="10 9 9 9 8 9"></polyline>
            </svg>
            <h4 style='margin: 0; color: #475569;'>Coaching Themes</h4>
        </div>
    """, unsafe_allow_html=True)
    theme_option = st.radio("Theme source:", ["Pre-loaded", "Custom", "Both"])
    
    if theme_option == "Custom":
        custom_themes = st.text_area("Enter themes (one per line):", height=150)
        coaching_themes = [t.strip() for t in custom_themes.split('\n') if t.strip()]
    elif theme_option == "Both":
        custom_themes = st.text_area("Add custom themes:", height=100)
        additional = [t.strip() for t in custom_themes.split('\n') if t.strip()]
        coaching_themes = DEFAULT_THEMES + additional
    else:
        coaching_themes = DEFAULT_THEMES
    
    st.session_state.coaching_themes = coaching_themes
    st.session_state.analysis_model = analysis_model if st.session_state.llm_provider == "openrouter" else analysis_model
    
    st.caption(f"{len(coaching_themes)} themes active")

# Main content
st.markdown("""
    <div style='text-align: center; padding: 40px 20px; background: linear-gradient(135deg, #f8fafc 0%, #e2e8f0 100%); border-radius: 20px; margin: 20px 0; box-shadow: 0 4px 6px rgba(0,0,0,0.05);'>
        <div style='display: flex; justify-content: center; margin-bottom: 20px;'>
            <svg width="56" height="56" viewBox="0 0 24 24" fill="none" stroke="#0ea5e9" stroke-width="2" stroke-linecap="round" stroke-linejoin="round">
                <circle cx="12" cy="12" r="10"></circle>
                <path d="M12 6v6l4 2"></path>
            </svg>
        </div>
        <h1 style='font-size: 3rem; font-weight: 700; color: #1e293b; margin: 0 0 10px 0;'>QA Coaching Intelligence</h1>
        <p style='font-size: 1.2rem; color: #64748b; margin: 0;'>Transform Every Call into Coaching Excellence</p>
    </div>
""", unsafe_allow_html=True)

# Tabs
tab1, tab2, tab3, tab4 = st.tabs(["📤 Upload & Process", "📊 Dashboard", "💬 Q&A Chat", "💾 Export & Session"])

with tab1:
    st.markdown("<div style='background: rgba(255,255,255,0.95); padding: 40px; border-radius: 20px; margin: 20px 0;'>", unsafe_allow_html=True)
    
    st.markdown("""
        <div style='display: flex; align-items: center; gap: 10px; margin-bottom: 15px;'>
            <svg width="24" height="24" viewBox="0 0 24 24" fill="none" stroke="#0ea5e9" stroke-width="2" stroke-linecap="round" stroke-linejoin="round">
                <path d="M21 15v4a2 2 0 0 1-2 2H5a2 2 0 0 1-2-2v-4"></path>
                <polyline points="17 8 12 3 7 8"></polyline>
                <line x1="12" y1="3" x2="12" y2="15"></line>
            </svg>
            <h3 style='margin: 0; color: #1e293b;'>Step 1: Upload Files</h3>
        </div>
    """, unsafe_allow_html=True)
    uploaded_files = st.file_uploader(
        "Supported: CSV, XLSX, XLS, TXT, Parquet",
        type=['csv', 'xlsx', 'xls', 'txt', 'parquet'],
        accept_multiple_files=True
    )
    
    if uploaded_files:
        st.success(f"✅ {len(uploaded_files)} file(s) uploaded")
        
        if st.button("🔄 Load & Convert to Parquet", use_container_width=True):
            with st.spinner("Loading and converting files..."):
                all_data = []
                
                for uploaded_file in uploaded_files:
                    df = load_file_to_dataframe(uploaded_file)
                    if df is not None:
                        all_data.append(df)
                
                if all_data:
                    combined_df = pd.concat(all_data, ignore_index=True)
                    
                    # Convert to parquet
                    parquet_bytes = convert_to_parquet(combined_df, 'transcripts.parquet')
                    st.session_state.transcripts_parquet = parquet_bytes
                    st.session_state.raw_df = combined_df
                    
                    # Load into DuckDB
                    conn = st.session_state.duckdb_conn
                    conn.execute("DROP TABLE IF EXISTS transcripts")
                    conn.execute("CREATE TABLE transcripts AS SELECT * FROM combined_df")
                    
                    st.success(f"✅ Loaded {len(combined_df):,} rows | Size: {len(parquet_bytes) / 1024 / 1024:.2f} MB (Parquet)")
                    
                    st.session_state.data_loaded = True
    
    if st.session_state.get('data_loaded'):
        st.markdown("---")
        st.markdown("""
            <div style='display: flex; align-items: center; gap: 10px; margin-bottom: 15px;'>
                <svg width="24" height="24" viewBox="0 0 24 24" fill="none" stroke="#0ea5e9" stroke-width="2" stroke-linecap="round" stroke-linejoin="round">
                    <path d="M21 16V8a2 2 0 0 0-1-1.73l-7-4a2 2 0 0 0-2 0l-7 4A2 2 0 0 0 3 8v8a2 2 0 0 0 1 1.73l7 4a2 2 0 0 0 2 0l7-4A2 2 0 0 0 21 16z"></path>
                    <polyline points="3.27 6.96 12 12.01 20.73 6.96"></polyline>
                    <line x1="12" y1="22.08" x2="12" y2="12"></line>
                </svg>
                <h3 style='margin: 0; color: #1e293b;'>Step 2: Map Columns</h3>
            </div>
        """, unsafe_allow_html=True)
        
        df = st.session_state.raw_df
        available_columns = list(df.columns)
        
        col1, col2, col3 = st.columns(3)
        
        with col1:
            st.markdown("""
                <div style='display: flex; align-items: center; gap: 8px; margin: 10px 0;'>
                    <svg width="18" height="18" viewBox="0 0 24 24" fill="none" stroke="#dc2626" stroke-width="2" stroke-linecap="round" stroke-linejoin="round">
                        <circle cx="12" cy="12" r="10"></circle>
                        <line x1="12" y1="8" x2="12" y2="12"></line>
                        <line x1="12" y1="16" x2="12.01" y2="16"></line>
                    </svg>
                    <h4 style='margin: 0; color: #475569;'>Required Fields</h4>
                </div>
            """, unsafe_allow_html=True)
            call_id_col = st.selectbox("Call ID column:", [""] + available_columns, key="call_id_col")
            agent_col = st.selectbox("Agent column:", [""] + available_columns, key="agent_col")
            transcript_col = st.selectbox("Transcript column:", [""] + available_columns, key="transcript_col")
        
        with col2:
            st.markdown("""
                <div style='display: flex; align-items: center; gap: 8px; margin: 10px 0;'>
                    <svg width="18" height="18" viewBox="0 0 24 24" fill="none" stroke="#f59e0b" stroke-width="2" stroke-linecap="round" stroke-linejoin="round">
                        <circle cx="12" cy="12" r="10"></circle>
                        <line x1="12" y1="16" x2="12" y2="12"></line>
                        <line x1="12" y1="8" x2="12.01" y2="8"></line>
                    </svg>
                    <h4 style='margin: 0; color: #475569;'>Optional Fields</h4>
                </div>
            """, unsafe_allow_html=True)
            sentiment_col = st.selectbox("Sentiment Score:", ["None"] + available_columns, key="sentiment_col")
            timestamp_col = st.selectbox("Timestamp:", ["None"] + available_columns, key="timestamp_col")
            duration_col = st.selectbox("Call Duration:", ["None"] + available_columns, key="duration_col")
        
        with col3:
            st.markdown("""
                <div style='display: flex; align-items: center; gap: 8px; margin: 10px 0;'>
                    <svg width="18" height="18" viewBox="0 0 24 24" fill="none" stroke="#10b981" stroke-width="2" stroke-linecap="round" stroke-linejoin="round">
                        <line x1="18" y1="20" x2="18" y2="10"></line>
                        <line x1="12" y1="20" x2="12" y2="4"></line>
                        <line x1="6" y1="20" x2="6" y2="14"></line>
                    </svg>
                    <h4 style='margin: 0; color: #475569;'>Additional Metrics</h4>
                </div>
            """, unsafe_allow_html=True)
            custom_cols = st.multiselect(
                "Other columns to include:",
                [c for c in available_columns if c not in [call_id_col, agent_col, transcript_col, sentiment_col, timestamp_col, duration_col]],
                key="custom_cols"
            )
        
        # Validate required fields
        if call_id_col and agent_col and transcript_col:
            st.success("✅ Required columns mapped")
            
            if st.button("📊 Run Pre-Analysis (DuckDB)", use_container_width=True):
                with st.spinner("Running analytics..."):
                    conn = st.session_state.duckdb_conn
                    
                    # Store column mapping
                    st.session_state.column_mapping = {
                        'call_id': call_id_col,
                        'agent': agent_col,
                        'transcript': transcript_col,
                        'sentiment': sentiment_col if sentiment_col != "None" else None,
                        'timestamp': timestamp_col if timestamp_col != "None" else None,
                        'duration': duration_col if duration_col != "None" else None,
                        'custom': custom_cols
                    }
                    
                    # Show progress
                    progress_text = st.empty()
                    progress_bar = st.progress(0.0)
                    
                    progress_text.text("📝 Parsing transcripts in parallel...")
                    progress_bar.progress(0.2)
                    
                    # Parse transcripts using parallel processing
                    from multiprocessing import cpu_count
                    from concurrent.futures import ProcessPoolExecutor, as_completed
                    
                    num_cores = cpu_count()
                    num_workers = max(1, num_cores - 1)
                    st.info(f"🚀 Using {num_workers} CPU cores for parallel processing")
                    
                    # Prepare data
                    chunk_data = []
                    for idx, row in df.iterrows():
                        call_id = row[call_id_col]
                        agent_name = row[agent_col]
                        transcript_text = row[transcript_col]
                        sentiment = None
                        if sentiment_col and sentiment_col != "None":
                            sentiment = row.get(sentiment_col)
                        
                        chunk_data.append((call_id, agent_name, transcript_text, sentiment))
                    
                    # Split into chunks
                    chunk_size = max(50, len(chunk_data) // (num_workers * 8))
                    chunks = [chunk_data[i:i + chunk_size] for i in range(0, len(chunk_data), chunk_size)]
                    total_chunks = len(chunks)
                    
                    progress_text.text(f"📝 Processing {len(chunk_data):,} transcripts in {total_chunks} chunks...")
                    
                    # Process in parallel with incremental progress
                    expanded_rows = []
                    completed_chunks = 0
                    
                    with ProcessPoolExecutor(max_workers=num_workers) as executor:
                        future_to_chunk = {executor.submit(parse_transcript_chunk, chunk): i for i, chunk in enumerate(chunks)}
                        
                        for future in as_completed(future_to_chunk):
                            try:
                                chunk_result = future.result(timeout=30)
                                expanded_rows.extend(chunk_result)
                                completed_chunks += 1
                                
                                # Update progress
                                progress_pct = 0.2 + (0.4 * completed_chunks / total_chunks)
                                progress_bar.progress(progress_pct)
                                progress_text.text(f"📝 Parsed {completed_chunks}/{total_chunks} chunks ({len(expanded_rows):,} turns so far)...")
                                
                            except Exception as e:
                                st.warning(f"⚠️ Chunk processing error: {str(e)}")
                                completed_chunks += 1
                                continue
                    
                    progress_bar.progress(0.6)
                    progress_text.text("💾 Loading into DuckDB...")
                    
                    if not expanded_rows:
                        progress_text.empty()
                        progress_bar.empty()
                        st.error("❌ No transcripts could be parsed. Please check your data format.")
                        st.info("Expected formats:\n- `[12:30:08 AGENT]: message`\n- `2025-02-07 13:17:57 +0000 Agent: message | 2025-02-07 13:18:01 +0000 Customer: response`")
                        st.stop()
                    
                    expanded_df = pd.DataFrame(expanded_rows)
                    
                    # Store in session state first
                    st.session_state.processed_df = expanded_df
                    
                    # Reload into DuckDB using direct DataFrame reference
                    conn.execute("DROP TABLE IF EXISTS transcripts")
                    conn.execute("CREATE TABLE transcripts AS SELECT * FROM expanded_df")
                    
                    progress_bar.progress(0.8)
                    progress_text.text("📊 Running analytics...")
                    
                    # Run DuckDB analytics
                    analytics = {}
                    
                    # 1. Call volumes
                    analytics['total_calls'] = conn.execute("SELECT COUNT(DISTINCT call_id) as count FROM transcripts").fetchone()[0]
                    analytics['total_agents'] = conn.execute("SELECT COUNT(DISTINCT agent) as count FROM transcripts").fetchone()[0]
                    
                    # 2. Per-agent stats
                    agent_stats = conn.execute("""
                        SELECT 
                            agent,
                            COUNT(DISTINCT call_id) as total_calls,
                            COUNT(*) as total_messages,
                            SUM(CASE WHEN speaker = 'agent' THEN 1 ELSE 0 END) as agent_messages,
                            SUM(CASE WHEN speaker = 'customer' THEN 1 ELSE 0 END) as customer_messages,
                            AVG(LENGTH(message)) as avg_message_length
                        FROM transcripts
                        GROUP BY agent
                        ORDER BY total_calls DESC
                    """).fetchdf()
                    analytics['agent_stats'] = agent_stats
                    
                    # 3. Sentiment analysis (if available)
                    if sentiment_col != "None":
                        sentiment_stats = conn.execute("""
                            SELECT 
                                agent,
                                AVG(sentiment_score) as avg_sentiment,
                                MIN(sentiment_score) as min_sentiment,
                                MAX(sentiment_score) as max_sentiment,
                                COUNT(CASE WHEN sentiment_score < 0.5 THEN 1 END) as low_sentiment_calls
                            FROM (
                                SELECT DISTINCT call_id, agent, sentiment_score 
                                FROM transcripts 
                                WHERE sentiment_score IS NOT NULL
                            )
                            GROUP BY agent
                        """).fetchdf()
                        analytics['sentiment_stats'] = sentiment_stats
                    
                    # 4. Message flow analysis
                    flow_stats = conn.execute("""
                        SELECT 
                            agent,
                            AVG(turns_per_call) as avg_turns,
                            AVG(agent_response_ratio) as avg_response_ratio
                        FROM (
                            SELECT 
                                call_id,
                                agent,
                                COUNT(*) as turns_per_call,
                                SUM(CASE WHEN speaker = 'agent' THEN 1 ELSE 0 END) * 1.0 / COUNT(*) as agent_response_ratio
                            FROM transcripts
                            GROUP BY call_id, agent
                        )
                        GROUP BY agent
                    """).fetchdf()
                    analytics['flow_stats'] = flow_stats
                    
                    progress_bar.progress(1.0)
                    progress_text.text("✅ Complete!")
                    
                    st.session_state.pre_analytics = analytics
                    st.session_state.pre_analysis_done = True
                    
                    # Clear progress indicators
                    import time
                    time.sleep(0.5)
                    progress_text.empty()
                    progress_bar.empty()
                    
                    st.success(f"✅ Pre-analysis complete! Processed {len(expanded_rows):,} message turns from {len(df):,} calls using {num_workers} CPU cores.")
                    st.rerun()
        else:
            st.warning("⚠️ Please map all required columns (Call ID, Agent, Transcript)")
    
    st.markdown("</div>", unsafe_allow_html=True)

with tab2:
    if st.session_state.get('pre_analysis_done'):
        analytics = st.session_state.pre_analytics
        
        st.markdown("""
            <div style='display: flex; align-items: center; gap: 10px; margin-bottom: 15px;'>
                <svg width="24" height="24" viewBox="0 0 24 24" fill="none" stroke="#0ea5e9" stroke-width="2" stroke-linecap="round" stroke-linejoin="round">
                    <line x1="18" y1="20" x2="18" y2="10"></line>
                    <line x1="12" y1="20" x2="12" y2="4"></line>
                    <line x1="6" y1="20" x2="6" y2="14"></line>
                </svg>
                <h3 style='margin: 0; color: #1e293b;'>Pre-Analysis Dashboard (DuckDB)</h3>
            </div>
        """, unsafe_allow_html=True)
        
        # Key metrics
        col1, col2, col3, col4 = st.columns(4)
        with col1:
            st.metric("Total Calls", f"{analytics['total_calls']:,}")
        with col2:
            st.metric("Total Agents", analytics['total_agents'])
        with col3:
            if 'sentiment_stats' in analytics:
                avg_sentiment = analytics['sentiment_stats']['avg_sentiment'].mean()
                st.metric("Avg Sentiment", f"{avg_sentiment:.2f}")
            else:
                st.metric("Avg Sentiment", "N/A")
        with col4:
            total_messages = analytics['agent_stats']['total_messages'].sum()
            st.metric("Total Messages", f"{total_messages:,}")
        
        st.markdown("---")
        
        # Agent performance table
        st.markdown("""
            <div style='display: flex; align-items: center; gap: 10px; margin: 15px 0;'>
                <svg width="22" height="22" viewBox="0 0 24 24" fill="none" stroke="#64748b" stroke-width="2" stroke-linecap="round" stroke-linejoin="round">
                    <path d="M17 21v-2a4 4 0 0 0-4-4H5a4 4 0 0 0-4 4v2"></path>
                    <circle cx="9" cy="7" r="4"></circle>
                    <path d="M23 21v-2a4 4 0 0 0-3-3.87"></path>
                    <path d="M16 3.13a4 4 0 0 1 0 7.75"></path>
                </svg>
                <h4 style='margin: 0; color: #1e293b;'>Agent Statistics</h4>
            </div>
        """, unsafe_allow_html=True)
        
        # Merge stats
        display_df = analytics['agent_stats'].copy()
        
        if 'sentiment_stats' in analytics:
            display_df = display_df.merge(
                analytics['sentiment_stats'][['agent', 'avg_sentiment', 'low_sentiment_calls']],
                on='agent',
                how='left'
            )
        
        if 'flow_stats' in analytics:
            display_df = display_df.merge(
                analytics['flow_stats'][['agent', 'avg_turns', 'avg_response_ratio']],
                on='agent',
                how='left'
            )
        
        # Format for display
        display_df['avg_message_length'] = display_df['avg_message_length'].round(1)
        if 'avg_sentiment' in display_df.columns:
            display_df['avg_sentiment'] = display_df['avg_sentiment'].round(2)
        if 'avg_turns' in display_df.columns:
            display_df['avg_turns'] = display_df['avg_turns'].round(1)
        if 'avg_response_ratio' in display_df.columns:
            display_df['avg_response_ratio'] = (display_df['avg_response_ratio'] * 100).round(1)
        
        st.dataframe(display_df, use_container_width=True, height=400)
        
        st.markdown("---")
        
        # Coaching insights section
        if not st.session_state.get('processed'):
            st.markdown("""
                <div style='display: flex; align-items: center; gap: 10px; margin-bottom: 15px;'>
                    <svg width="24" height="24" viewBox="0 0 24 24" fill="none" stroke="#0ea5e9" stroke-width="2" stroke-linecap="round" stroke-linejoin="round">
                        <circle cx="12" cy="12" r="10"></circle>
                        <circle cx="12" cy="12" r="3"></circle>
                    </svg>
                    <h3 style='margin: 0; color: #1e293b;'>Generate Coaching Insights</h3>
                </div>
            """, unsafe_allow_html=True)
            st.info("📌 Pre-analysis complete! Now optionally generate AI-powered coaching themes.")
            
            # Validation
            if st.session_state.llm_provider == "openrouter" and not st.session_state.get('openrouter_api_key'):
                st.error("❌ Please provide OpenRouter API key in sidebar")
            else:
                if st.button("Generate Coaching Themes (LLM)", use_container_width=True, type="primary"):
                    with st.spinner("Generating coaching insights..."):
                        # Get agents data
                        df = st.session_state.processed_df
                        
                        # Calculate file hash for session management
                        file_hash = calculate_file_hash(df)
                        st.session_state.file_hash = file_hash
                        
                        # Check for existing session
                        existing_session = load_latest_session(file_hash)
                        previously_analyzed = set(existing_session['agent_names']) if existing_session else set()
                        
                        # Filter agents with 10+ calls
                        agent_call_counts = df.groupby('agent')['call_id'].nunique()
                        eligible_agents = agent_call_counts[agent_call_counts >= 10].index.tolist()
                        excluded_agents = agent_call_counts[agent_call_counts < 10]
                        
                        # Remove already analyzed agents
                        new_agents = [a for a in eligible_agents if a not in previously_analyzed]
                        
                        # Filter DataFrame to eligible agents only
                        df_filtered = df[df['agent'].isin(new_agents)]
                        
                        # Show filtering info
                        info_msg = f"**Analysis Scope:**\n"
                        info_msg += f"- Analyzing {len(new_agents)} new agents with 10+ calls\n"
                        info_msg += f"- Excluded {len(excluded_agents)} agents with <10 calls\n"
                        if previously_analyzed:
                            info_msg += f"- Resuming from previous session ({len(previously_analyzed)} agents already analyzed)\n"
                        info_msg += f"- Prioritizing low sentiment calls for coaching focus"
                        st.info(info_msg)
                        
                        if not new_agents:
                            st.success("All eligible agents already analyzed! Loading previous results...")
                            st.session_state.coaching_insights = existing_session['insights']
                            st.session_state.processed = True
                            st.rerun()
                        
                        # Group by agent and prioritize low sentiment calls
                        agents_data = []
                        for agent in new_agents:
                            agent_df = df_filtered[df_filtered['agent'] == agent]
                            
                            # Get unique calls for this agent
                            call_ids = agent_df['call_id'].unique()
                            
                            # Calculate sentiment per call (if available)
                            if 'sentiment_score' in agent_df.columns:
                                call_sentiments = []
                                for call_id in call_ids:
                                    call_data = agent_df[agent_df['call_id'] == call_id]
                                    avg_sentiment = call_data['sentiment_score'].mean()
                                    if pd.notna(avg_sentiment):
                                        call_sentiments.append((call_id, avg_sentiment))
                                
                                # Sort by sentiment (lowest first) and take top 5
                                if call_sentiments:
                                    call_sentiments.sort(key=lambda x: x[1])
                                    selected_calls = [c[0] for c in call_sentiments[:5]]
                                else:
                                    selected_calls = call_ids[:5]
                            else:
                                selected_calls = call_ids[:5]
                            
                            # Get data for selected calls
                            selected_data = agent_df[agent_df['call_id'].isin(selected_calls)]
                            agents_data.append((agent, selected_data))
                        
                        total_agents = len(agents_data)
                        themes = st.session_state.get('coaching_themes', DEFAULT_THEMES)
                        
                        # Batch configuration
                        BATCH_SIZE = 20
                        agent_batches = [agents_data[i:i + BATCH_SIZE] for i in range(0, len(agents_data), BATCH_SIZE)]
                        total_batches = len(agent_batches)
                        
                        st.info(f"Processing {total_agents} agents in {total_batches} batch(es) of up to {BATCH_SIZE} agents each")
                        
                        # Initialize or load insights
                        all_insights = existing_session['insights'].copy() if existing_session else {}
                        
                        # Progress tracking
                        overall_progress = st.progress(0.0)
                        batch_status = st.empty()
                        
                        start_time = time.time()
                        
                        # Process each batch
                        for batch_num, agent_batch in enumerate(agent_batches, 1):
                            batch_status.write(f"**Batch {batch_num}/{total_batches}** - Processing {len(agent_batch)} agents...")
                            
                            # Run parallel processing for this batch
                            async def run_batch():
                                return await process_all_agents_parallel(
                                    agent_batch,
                                    themes,
                                    st.session_state.get('analysis_model', 'deepseek/deepseek-chat:free'),
                                    st.session_state.llm_provider,
                                    st.session_state.get('openrouter_api_key'),
                                    st.session_state.get('local_llm_url'),
                                    max_concurrent=st.session_state.get('max_concurrent', 10),
                                    calls_per_minute=st.session_state.get('calls_per_minute', 50)
                                )
                            
                            import nest_asyncio
                            nest_asyncio.apply()
                            
                            loop = asyncio.new_event_loop()
                            asyncio.set_event_loop(loop)
                            
                            try:
                                batch_insights = loop.run_until_complete(run_batch())
                                
                                # Merge results
                                all_insights.update(batch_insights)
                                
                                # Save session after each batch
                                save_session(
                                    file_hash=file_hash,
                                    insights=all_insights,
                                    filter_criteria={'min_calls': 10, 'sentiment_priority': True},
                                    model_used=st.session_state.get('analysis_model', 'deepseek/deepseek-chat:free')
                                )
                                
                                batch_status.success(f"Batch {batch_num}/{total_batches} complete: {len(batch_insights)} agents processed")
                                overall_progress.progress(batch_num / total_batches)
                                
                            except Exception as e:
                                batch_status.error(f"Batch {batch_num} error: {str(e)}")
                                continue
                        
                        elapsed = time.time() - start_time
                        overall_progress.progress(1.0)
                        
                        if not all_insights or len(all_insights) == 0:
                            st.error("No insights generated. The LLM may have failed.")
                            
                            # Retry button
                            if st.button("Retry with Selected Model", use_container_width=True, type="primary"):
                                st.session_state.processed = False
                                st.session_state.coaching_insights = {}
                                st.rerun()
                        else:
                            batch_status.success(f"All batches complete! Processed {len(all_insights)} agents in {elapsed:.1f}s")
                            
                            # Save insights to DuckDB for caching and chat context
                            conn = st.session_state.duckdb_conn
                            cache_rows = []
                            for agent_name, agent_data in all_insights.items():
                                for theme in agent_data.get('coaching_themes', []):
                                    cache_rows.append({
                                        'agent': agent_name,
                                        'theme': theme.get('theme', ''),
                                        'priority': theme.get('priority', 'low'),
                                        'frequency': theme.get('frequency', 1),
                                        'examples': str(theme.get('examples', [])),
                                        'recommendation': theme.get('recommendation', ''),
                                        'processed_at': datetime.now().isoformat(),
                                        'model_used': st.session_state.get('analysis_model', 'unknown')
                                    })
                            
                            if cache_rows:
                                cache_df = pd.DataFrame(cache_rows)
                                conn.execute("DROP TABLE IF EXISTS coaching_cache")
                                conn.execute("CREATE TABLE coaching_cache AS SELECT * FROM cache_df")
                                st.success(f"Cached {len(cache_rows)} coaching insights for future queries")
                            
                            # Generate analytics context for chat
                            st.session_state.analytics_context = generate_analytics_context(all_insights, df)
                            
                            st.session_state.coaching_insights = all_insights
                            st.session_state.processed = True
                            time.sleep(1)
                            st.rerun()
        
        # Show coaching insights if available
        if st.session_state.get('processed'):
            st.markdown("---")
            st.markdown("""
                <div style='display: flex; align-items: center; gap: 10px; margin-bottom: 20px;'>
                    <svg width="24" height="24" viewBox="0 0 24 24" fill="none" stroke="#10b981" stroke-width="2" stroke-linecap="round" stroke-linejoin="round">
                        <polyline points="22 12 18 12 15 21 9 3 6 12 2 12"></polyline>
                    </svg>
                    <h3 style='margin: 0; color: #1e293b;'>AI-Powered Coaching Insights</h3>
                </div>
            """, unsafe_allow_html=True)
            
            insights = st.session_state.coaching_insights
            df = st.session_state.processed_df
            
            # Generate HTML report
            html_report = generate_html_report(insights, df)
            st.components.v1.html(html_report, height=2000, scrolling=True)
            
            st.session_state.html_report = html_report
    
    else:
        st.info("👆 Upload files and run pre-analysis first!")

with tab3:
    if st.session_state.processed:
        st.markdown("""
            <div style='display: flex; align-items: center; gap: 10px; margin-bottom: 15px;'>
                <svg width="24" height="24" viewBox="0 0 24 24" fill="none" stroke="#0ea5e9" stroke-width="2" stroke-linecap="round" stroke-linejoin="round">
                    <path d="M21 15a2 2 0 0 1-2 2H7l-4 4V5a2 2 0 0 1 2-2h14a2 2 0 0 1 2 2z"></path>
                </svg>
                <h3 style='margin: 0; color: #1e293b;'>Ask Questions About Your Data</h3>
            </div>
        """, unsafe_allow_html=True)
        
        # Chat model selector
        with st.expander("⚙️ Chat Settings"):
            chat_model = st.selectbox(
                "Chat model:",
                options=list(MODELS.keys()),
                format_func=lambda x: f"[{MODELS[x]['rating']}] {MODELS[x]['name']}",
                index=list(MODELS.keys()).index("mistralai/mistral-nemo:free")
            )
        
        # Display chat history
        for msg in st.session_state.chat_history:
            if msg['role'] == 'user':
                st.markdown(f"<div class='chat-message user-message'><strong>You:</strong> {msg['content']}</div>", unsafe_allow_html=True)
            else:
                st.markdown(f"<div class='chat-message assistant-message'><strong>Assistant:</strong> {msg['content']}</div>", unsafe_allow_html=True)
        
        # Chat input
        user_question = st.text_input("Ask a question:", key="chat_input")
        
        col1, col2 = st.columns([1, 5])
        with col1:
            send_btn = st.button("Send", use_container_width=True)
        with col2:
            clear_btn = st.button("Clear History", use_container_width=True)
        
        if clear_btn:
            st.session_state.chat_history = []
            st.rerun()
        
        if send_btn and user_question:
            # Add to history
            st.session_state.chat_history.append({"role": "user", "content": user_question})
            
            # Prepare context
            insights = st.session_state.coaching_insights
            df = st.session_state.processed_df
            analytics_context = st.session_state.get('analytics_context', '')
            
            # Simple query routing
            question_lower = user_question.lower()
            
            # Check coaching cache first
            try:
                conn = st.session_state.duckdb_conn
                has_coaching = conn.execute("SELECT COUNT(*) FROM coaching_cache").fetchone()[0] > 0
                
                # Quick SQL queries for specific patterns
                if has_coaching and any(kw in question_lower for kw in ['coaching', 'theme', 'improve', 'recommendation']):
                    # Coaching-specific queries
                    if 'top' in question_lower and 'theme' in question_lower:
                        result = conn.execute("""
                            SELECT theme, COUNT(*) as count 
                            FROM coaching_cache 
                            GROUP BY theme 
                            ORDER BY count DESC 
                            LIMIT 5
                        """).fetchdf()
                        themes_list = "\n".join([f"- {row['theme']}: {row['count']} agents" for _, row in result.iterrows()])
                        answer = f"**Top coaching themes:**\n{themes_list}"
                        st.session_state.chat_history.append({"role": "assistant", "content": answer})
                        st.rerun()
                    
                    # Check if agent name is in question
                    for agent in insights.keys():
                        if agent.lower() in question_lower:
                            result = conn.execute(f"""
                                SELECT theme, priority, recommendation 
                                FROM coaching_cache 
                                WHERE agent = '{agent}' 
                                ORDER BY 
                                    CASE priority 
                                        WHEN 'high' THEN 1 
                                        WHEN 'medium' THEN 2 
                                        ELSE 3 
                                    END
                                LIMIT 3
                            """).fetchdf()
                            themes = "\n".join([f"- **{row['theme']}** ({row['priority']} priority): {row['recommendation']}" 
                                               for _, row in result.iterrows()])
                            answer = f"**Coaching needs for {agent}:**\n{themes}"
                            st.session_state.chat_history.append({"role": "assistant", "content": answer})
                            st.rerun()
            except:
                pass
            
            # Check if it's a SQL-like question
            if any(kw in question_lower for kw in ['how many', 'count', 'average', 'total', 'list all']):
                # Try to answer with DuckDB
                try:
                    conn = st.session_state.duckdb_conn
                    
                    if 'how many calls' in question_lower:
                        result = conn.execute("SELECT COUNT(DISTINCT call_id) as count FROM transcripts").fetchone()
                        answer = f"There are {result[0]} calls in the dataset."
                    elif 'how many agents' in question_lower:
                        result = conn.execute("SELECT COUNT(DISTINCT agent) as count FROM transcripts WHERE agent IS NOT NULL").fetchone()
                        answer = f"There are {result[0]} agents in the dataset."
                    elif 'average sentiment' in question_lower:
                        if 'sentiment_score' in df.columns:
                            result = conn.execute("SELECT AVG(sentiment_score) as avg FROM transcripts").fetchone()
                            answer = f"The average sentiment score is {result[0]:.2f}."
                        else:
                            answer = "Sentiment data not available in the transcripts."
                    else:
                        answer = "I can answer questions about call counts, agent counts, and sentiment averages. Try rephrasing!"
                        
                    st.session_state.chat_history.append({"role": "assistant", "content": answer})
                    st.rerun()
                    
                except Exception as e:
                    answer = f"I encountered an error: {str(e)}"
                    st.session_state.chat_history.append({"role": "assistant", "content": answer})
                    st.rerun()
            
            else:
                # Use LLM for complex analytical questions with full analytics context
                enhanced_context = f"""{analytics_context}

Detailed Coaching Insights:
"""
                for agent, data in list(insights.items())[:20]:  # Limit to prevent token overflow
                    enhanced_context += f"\n{agent}: {len(data.get('coaching_themes', []))} themes\n"
                    for theme in data.get('coaching_themes', [])[:2]:
                        enhanced_context += f"  - {theme.get('theme', '')} ({theme.get('priority', '')}): {theme.get('recommendation', '')[:100]}...\n"
                
                messages = [
                    {"role": "system", "content": "You are an expert QA coaching analyst with access to comprehensive coaching data and analytics. Answer questions with specific insights, data-driven recommendations, and cite agent names when relevant. Use the analytics context to provide accurate statistics and trends."},
                    {"role": "user", "content": f"Analytics Context:\n{enhanced_context}\n\nQuestion: {user_question}\n\nProvide a detailed, data-driven answer based on the analytics context above."}
                ]
                
                with st.spinner("Analyzing with AI..."):
                    response = call_llm(chat_model, messages, temperature=0.5, is_json=False)
                
                if response and 'choices' in response:
                    answer = response['choices'][0]['message']['content']
                    st.session_state.chat_history.append({"role": "assistant", "content": answer})
                    st.rerun()
                else:
                    st.error("Failed to get response from chat model")
    else:
        st.info("👆 Process transcripts first to enable Q&A!")

with tab4:
    if st.session_state.processed:
        st.markdown("""
            <div style='display: flex; align-items: center; gap: 10px; margin-bottom: 15px;'>
                <svg width="24" height="24" viewBox="0 0 24 24" fill="none" stroke="#0ea5e9" stroke-width="2" stroke-linecap="round" stroke-linejoin="round">
                    <path d="M21 15v4a2 2 0 0 1-2 2H5a2 2 0 0 1-2-2v-4"></path>
                    <polyline points="7 10 12 15 17 10"></polyline>
                    <line x1="12" y1="15" x2="12" y2="3"></line>
                </svg>
                <h3 style='margin: 0; color: #1e293b;'>Download Reports</h3>
            </div>
        """, unsafe_allow_html=True)
        
        col1, col2, col3 = st.columns(3)
        
        with col1:
            # HTML Report
            if 'html_report' in st.session_state:
                st.download_button(
                    "📊 HTML Report",
                    data=st.session_state.html_report,
                    file_name=f"coaching_report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.html",
                    mime="text/html",
                    use_container_width=True
                )
        
        with col2:
            # CSV Export
            if st.button("📄 Export CSV", use_container_width=True):
                insights = st.session_state.coaching_insights
                rows = []
                for agent, data in insights.items():
                    for theme in data.get('coaching_themes', []):
                        rows.append({
                            'agent': agent,
                            'theme': theme.get('theme', ''),
                            'priority': theme.get('priority', ''),
                            'frequency': theme.get('frequency', 0),
                            'examples': ' | '.join(theme.get('examples', [])),
                            'recommendation': theme.get('recommendation', '')
                        })
                
                export_df = pd.DataFrame(rows)
                csv = export_df.to_csv(index=False)
                
                st.download_button(
                    "Download CSV",
                    data=csv,
                    file_name=f"coaching_data_{datetime.now().strftime('%Y%m%d_%H%M%S')}.csv",
                    mime="text/csv",
                    use_container_width=True
                )
        
        with col3:
            # Excel Export
            if st.button("📗 Export Excel", use_container_width=True):
                insights = st.session_state.coaching_insights
                
                # Create Excel file
                output = io.BytesIO()
                with pd.ExcelWriter(output, engine='openpyxl') as writer:
                    # Summary sheet
                    summary_data = []
                    for agent, data in insights.items():
                        summary_data.append({
                            'Agent': agent,
                            'Calls': data.get('calls_analyzed', 0),
                            'Coaching Themes': len(data.get('coaching_themes', [])),
                            'High Priority': sum(1 for t in data.get('coaching_themes', []) if t.get('priority') == 'high')
                        })
                    pd.DataFrame(summary_data).to_excel(writer, sheet_name='Summary', index=False)
                    
                    # Detail sheet
                    rows = []
                    for agent, data in insights.items():
                        for theme in data.get('coaching_themes', []):
                            rows.append({
                                'Agent': agent,
                                'Theme': theme.get('theme', ''),
                                'Priority': theme.get('priority', ''),
                                'Frequency': theme.get('frequency', 0),
                                'Recommendation': theme.get('recommendation', '')
                            })
                    pd.DataFrame(rows).to_excel(writer, sheet_name='Coaching Details', index=False)
                
                excel_data = output.getvalue()
                
                st.download_button(
                    "Download Excel",
                    data=excel_data,
                    file_name=f"coaching_report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx",
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    use_container_width=True
                )
        
        # PowerPoint export
        col4, col5, col6 = st.columns(3)
        with col4:
            if st.button("📽️ Export PowerPoint", use_container_width=True):
                with st.spinner("Generating PowerPoint..."):
                    ppt_data = generate_powerpoint(st.session_state.coaching_insights, st.session_state.processed_df)
                    st.download_button(
                        "Download PowerPoint",
                        data=ppt_data,
                        file_name=f"coaching_report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        use_container_width=True
                    )
        
        st.markdown("---")
        st.markdown("""
            <div style='display: flex; align-items: center; gap: 10px; margin-bottom: 15px;'>
                <svg width="24" height="24" viewBox="0 0 24 24" fill="none" stroke="#0ea5e9" stroke-width="2" stroke-linecap="round" stroke-linejoin="round">
                    <path d="M19 21H5a2 2 0 0 1-2-2V5a2 2 0 0 1 2-2h11l5 5v11a2 2 0 0 1-2 2z"></path>
                    <polyline points="17 21 17 13 7 13 7 21"></polyline>
                    <polyline points="7 3 7 8 15 8"></polyline>
                </svg>
                <h3 style='margin: 0; color: #1e293b;'>Session Management</h3>
            </div>
        """, unsafe_allow_html=True)
        
        col1, col2 = st.columns(2)
        
        with col1:
            if st.button("💾 Save Session", use_container_width=True):
                # Save all session data to parquet
                session_data = {
                    'transcripts': st.session_state.processed_df,
                    'insights': pd.DataFrame([
                        {'agent': agent, 'insights_json': json.dumps(data)}
                        for agent, data in st.session_state.coaching_insights.items()
                    ]),
                    'metadata': pd.DataFrame([{
                        'processed_at': datetime.now().isoformat(),
                        'model_used': analysis_model,
                        'total_calls': len(st.session_state.processed_df['call_id'].unique()),
                        'total_agents': len(st.session_state.coaching_insights)
                    }])
                }
                
                # Convert each to parquet
                output = io.BytesIO()
                
                # For now, save insights as JSON in parquet
                session_df = pd.DataFrame([{
                    'data_type': 'insights',
                    'content': json.dumps(st.session_state.coaching_insights)
                }])
                
                parquet_bytes = convert_to_parquet(session_df, 'session.parquet')
                
                st.download_button(
                    "📥 Download Session File",
                    data=parquet_bytes,
                    file_name=f"session_{datetime.now().strftime('%Y%m%d_%H%M%S')}.parquet",
                    mime="application/octet-stream",
                    use_container_width=True
                )
        
        with col2:
            session_file = st.file_uploader("📂 Load Session", type=['parquet'])
            if session_file:
                # Load session
                session_df = pd.read_parquet(session_file)
                if not session_df.empty:
                    content = json.loads(session_df.iloc[0]['content'])
                    st.session_state.coaching_insights = content
                    st.session_state.processed = True
                    st.success("✅ Session loaded!")
                    st.rerun()
    
    else:
        st.info("👆 Process transcripts first!")

# Footer
st.markdown("<br><br>", unsafe_allow_html=True)
st.markdown("<div style='text-align: center; color: white; opacity: 0.7; padding: 20px;'>", unsafe_allow_html=True)
st.markdown("QA Coaching Intelligence Platform | Developed by CE INNOVATIONS LAB 2025", unsafe_allow_html=True)
st.markdown("</div>", unsafe_allow_html=True)
